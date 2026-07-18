from flask import Flask, request, render_template, jsonify, send_file
import requests
import anthropic
import json
import os
import re
import io
import time
import threading
import uuid
import concurrent.futures
from datetime import datetime, timezone
from collections import Counter
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from bs4 import BeautifulSoup

load_dotenv()
SEARCHAPI_KEY     = os.getenv("SEARCHAPI_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
APIFY_API_KEY     = os.getenv("APIFY_API_KEY")
YOUTUBE_API_KEY   = os.getenv("YOUTUBE_API_KEY")

app = Flask(__name__)

TOTAL_SLIDES = 40

# ═══════════════════════════════════════════════════════════════
# MIDNIGHT PANDA BRAND DESIGN TOKENS
# ═══════════════════════════════════════════════════════════════

BG_DARK        = RGBColor(0x0A, 0x0A, 0x0A)
BG_LIGHT       = RGBColor(0xF5, 0xF3, 0xEF)
GOLD           = RGBColor(0xB8, 0x94, 0x5A)
CARD_LIGHT     = RGBColor(0xEC, 0xE9, 0xE2)
CARD_DARK      = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_DARK      = RGBColor(0x0A, 0x0A, 0x0A)
TEXT_LIGHT     = RGBColor(0xF5, 0xF3, 0xEF)
TEXT_GRAY      = RGBColor(0x66, 0x66, 0x66)
TEXT_GRAY_LT   = RGBColor(0xD0, 0xCF, 0xC9)
TEXT_FOOTER    = RGBColor(0x99, 0x99, 0x99)
PIE_TAN        = RGBColor(0xD9, 0xC4, 0xA0)
MOMENTUM_UP    = RGBColor(0x6F, 0xA8, 0x6B)
MOMENTUM_DOWN  = RGBColor(0xC1, 0x6E, 0x6E)

FONT_MONO      = "DM Mono"
FONT_MONO_MED  = "DM Mono Medium"
FONT_MONO_LT   = "DM Mono Light"
FONT_SERIF     = "Playfair Display Black"
FONT_STAT      = "Bebas Neue"


# ═══════════════════════════════════════════════════════════════
# STEP 1 — DISCOVER SOCIAL HANDLES
# ═══════════════════════════════════════════════════════════════

def clean_domain(url):
    url = url.strip().rstrip("/")
    if not url.startswith("http"):
        url = "https://" + url
    return url

def is_youtube_channel_id(handle):
    raw = handle.lstrip("@")
    return raw.startswith("UC") and len(raw) >= 20

def extract_social_from_html(html):
    soup = BeautifulSoup(html, "html.parser")
    handles = {}
    SKIP = {"p","explore","accounts","stories","reel","tv","share","sharer",
            "login","signup","intent","home","search"}
    for tag in soup.find_all("a", href=True):
        href = tag["href"].lower()

        if "instagram.com" in href and "instagram" not in handles:
            m = re.search(r'instagram\.com/([A-Za-z0-9_.]+)/?', href, re.I)
            if m and m.group(1).lower() not in SKIP:
                handles["instagram"] = m.group(1)

        if "facebook.com" in href and "facebook" not in handles:
            if "profile.php" not in href:
                m = re.search(r'facebook\.com/([A-Za-z0-9_.]+)/?', href, re.I)
                if m and m.group(1).lower() not in SKIP:
                    handles["facebook"] = m.group(1)

        if "linkedin.com/company" in href and "linkedin" not in handles:
            m = re.search(r'linkedin\.com/company/([A-Za-z0-9_.-]+)/?', href, re.I)
            if m:
                handles["linkedin"] = m.group(1)

        if "youtube.com" in href and "youtube" not in handles:
            m = re.search(r'youtube\.com/@([A-Za-z0-9_.]+)/?', href, re.I)
            if m:
                handles["youtube"] = "@" + m.group(1)
            else:
                m = re.search(r'youtube\.com/(?:channel|c|user)/([A-Za-z0-9_.-]+)/?', href, re.I)
                if m:
                    handles["youtube"] = m.group(1)
    return handles

def find_handles_from_website(website_url):
    base = clean_domain(website_url)
    pages = [base, base+"/contact", base+"/about",
             base+"/contact-us", base+"/about-us"]
    headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) Chrome/120.0.0.0"}
    all_handles = {}
    for page in pages:
        try:
            r = requests.get(page, headers=headers, timeout=10)
            if r.status_code == 200:
                found = extract_social_from_html(r.text)
                all_handles.update(found)
                if len(all_handles) >= 4:
                    break
        except:
            continue
    return all_handles

def find_handles_from_searchapi(website_url, missing_platforms):
    domain = website_url.replace("https://","").replace("http://","").rstrip("/")
    brand  = re.sub(r'\.(com|in|io|co|net|org|app).*','', domain).lower()
    brand  = brand.replace("www.","")
    handles = {}
    SKIP = {"p","explore","accounts","stories","reel","tv","share",
            "sharer","login","signup","intent","home","search"}

    brand_words = re.split(r'[-_.]', brand)
    print(f"   Brand: '{brand}' | Words: {brand_words} | Domain: '{domain}'")

    queries = {
        "instagram": [
            f'{domain} instagram',
            f'"{brand}" instagram official',
            f'site:instagram.com "{brand}"',
        ],
        "facebook": [
            f'"{brand}" facebook page',
            f'{domain} facebook',
            f'site:facebook.com "{brand}"',
        ],
        "linkedin": [
            f'"{brand}" site:linkedin.com/company',
            f'{domain} linkedin',
            f'"{brand}" linkedin company followers',
        ],
        "youtube": [
            f'"{brand}" youtube channel',
            f'{domain} youtube',
        ],
    }

    for platform in missing_platforms:
        found = False
        for q in queries[platform]:
            if found:
                break
            try:
                print(f"   🔍 {platform}: {q}")
                r = requests.get(
                    "https://www.searchapi.io/api/v1/search",
                    params={"engine":"google","q":q,
                            "api_key":SEARCHAPI_KEY,"num":5}
                )
                if r.status_code != 200:
                    continue

                results = r.json().get("organic_results",[])
                for idx, result in enumerate(results):
                    link    = result.get("link","")
                    title   = result.get("title","").lower()
                    snippet = result.get("snippet","").lower()

                    def brand_match(text):
                        text_clean = text.replace("-","").replace("_","").replace(" ","").lower()
                        brand_clean = brand.replace("-","").replace("_","").replace(" ","")
                        if brand_clean in text_clean:
                            return True
                        for word in brand_words:
                            if len(word) > 3 and word in text_clean:
                                return True
                        return False

                    if platform == "instagram" and "instagram.com" in link:
                        m = re.search(r'instagram\.com/([A-Za-z0-9_.]+)/?', link, re.I)
                        if m:
                            handle = m.group(1)
                            if handle.lower() in SKIP:
                                continue
                            if brand_match(handle) or brand_match(title) or brand_match(snippet):
                                handles["instagram"] = handle
                                print(f"   ✅ Instagram: @{handle}")
                                found = True; break
                            elif idx == 0 and (domain.replace("www.","").split(".")[0] in handle.lower()
                                               or brand.split("-")[0] in handle.lower()):
                                handles["instagram"] = handle
                                print(f"   ✅ Instagram (domain match): @{handle}")
                                found = True; break
                            else:
                                print(f"   ⚠️ Instagram skipped: @{handle}")

                    elif platform == "facebook" and "facebook.com" in link:
                        if "profile.php" in link:
                            print(f"   ⚠️ Facebook skipped profile.php")
                            continue
                        m = re.search(r'facebook\.com/([A-Za-z0-9_.]+)/?', link, re.I)
                        if m:
                            handle = m.group(1)
                            if handle.lower() in SKIP:
                                continue
                            if brand_match(handle) or brand_match(title) or brand_match(snippet):
                                handles["facebook"] = handle
                                print(f"   ✅ Facebook: {handle}")
                                found = True; break
                            else:
                                print(f"   ⚠️ Facebook skipped wrong brand: {handle}")

                    elif platform == "linkedin" and "linkedin.com/company" in link:
                        m = re.search(r'linkedin\.com/company/([A-Za-z0-9_.-]+)/?', link, re.I)
                        if m:
                            handles["linkedin"] = m.group(1)
                            print(f"   ✅ LinkedIn: {m.group(1)}")
                            found = True; break

                    elif platform == "youtube" and "youtube.com" in link:
                        m = re.search(r'youtube\.com/@([A-Za-z0-9_.]+)/?', link, re.I)
                        if not m:
                            m = re.search(r'youtube\.com/(?:channel|c|user)/([A-Za-z0-9_.-]+)/?', link, re.I)
                        if m:
                            handle = m.group(1)
                            if brand_match(handle) or brand_match(title) or brand_match(snippet):
                                if is_youtube_channel_id(handle):
                                    handles["youtube"] = handle.lstrip("@")
                                else:
                                    handles["youtube"] = "@" + handle.lstrip("@")
                                print(f"   ✅ YouTube: {handles['youtube']}")
                                found = True; break
                            else:
                                print(f"   ⚠️ YouTube skipped: {handle}")

            except Exception as e:
                print(f"   ⚠️ SearchAPI failed for {platform}: {e}")

    return handles

def find_youtube_via_engine(brand):
    try:
        print(f"   🔍 YouTube engine search: {brand}")
        r = requests.get(
            "https://www.searchapi.io/api/v1/search",
            params={"engine":"youtube","q":brand,"api_key":SEARCHAPI_KEY}
        )
        if r.status_code != 200:
            return None
        brand_clean = brand.replace("-","").replace("_","").replace(" ","").lower()
        for video in r.json().get("videos", [])[:5]:
            ch = video.get("channel", {})
            ch_id    = ch.get("id","")
            ch_title = ch.get("title","").replace(" ","").replace("-","").lower()
            if ch_id and brand_clean in ch_title:
                print(f"   ✅ YouTube via engine: {ch.get('title')} ({ch_id})")
                return ch_id
    except Exception as e:
        print(f"   ⚠️ YouTube engine search failed: {e}")
    return None

def discover_all_handles(website_url):
    print(f"\n🔎 Finding all social handles for: {website_url}")
    handles = find_handles_from_website(website_url)
    print(f"   Found on website: {handles}")
    missing = [p for p in ["instagram","facebook","linkedin","youtube"]
               if p not in handles]
    if missing:
        print(f"   Searching for missing: {missing}")
        google_handles = find_handles_from_searchapi(website_url, missing)
        handles.update(google_handles)
        print(f"   Final handles: {handles}")

    if "youtube" not in handles:
        domain = website_url.replace("https://","").replace("http://","").rstrip("/")
        brand  = re.sub(r'\.(com|in|io|co|net|org|app).*','', domain).replace("www.","")
        yt_id = find_youtube_via_engine(brand)
        if yt_id:
            handles["youtube"] = yt_id
            print(f"   Final handles: {handles}")

    if not handles:
        raise ValueError(
            f"Could not find any social media accounts for {website_url}. "
            "Please check the website URL."
        )
    return handles


# ═══════════════════════════════════════════════════════════════
# STEP 2 — FETCH DATA FROM EACH PLATFORM
# ═══════════════════════════════════════════════════════════════

def _get_first(d, keys, default=None):
    for k in keys:
        v = d.get(k)
        if v not in (None, ""):
            return v
    return default

def _to_timestamp(val):
    """
    Convert various timestamp formats to a unix timestamp (epoch seconds).

    IMPORTANT — timezone handling:
    Instagram / Apify / SearchAPI all report post dates in UTC. The trailing
    "Z" on strings like "2026-04-01T13:30:41.000Z" means "Zulu time", i.e. UTC.
    We explicitly attach UTC tzinfo for naive date strings instead of letting
    Python guess based on wherever the server happens to be running — otherwise
    the exact same post produces a different timestamp on a Mac in IST vs a
    Railway server, which silently throws off posting frequency, day-of-week
    calculations, and the pinned-post date-sort below.
    """
    if val is None:
        return None
    try:
        if isinstance(val, (int, float)):
            v = float(val)
            # Millisecond-epoch detection: HarvestAPI (and some other
            # APIs) send epoch MILLISECONDS. Read as seconds, those put
            # posts in the year ~58509 and crash every date calculation
            # (confirmed live: "ValueError: year 58509 is out of range").
            # Any value above 1e11 seconds (~year 5138) can't be a real
            # post date in seconds, so it must be milliseconds.
            if v > 1e11:
                v = v / 1000.0
            return v
        s = str(val)
        if s.isdigit():
            v = float(s)
            if v > 1e11:
                v = v / 1000.0
            return v

        # Formats that already carry explicit timezone/offset info
        for fmt in ("%Y-%m-%dT%H:%M:%S.%f%z", "%Y-%m-%dT%H:%M:%S%z"):
            try:
                dt = datetime.strptime(s[:32], fmt)
                return dt.timestamp()
            except:
                continue

        # Formats with NO timezone marker — Instagram data here is always UTC
        for fmt in ("%Y-%m-%dT%H:%M:%S.%f", "%Y-%m-%dT%H:%M:%S",
                    "%Y-%m-%d %H:%M:%S", "%Y-%m-%d"):
            try:
                clean = s[:26].rstrip("Z")
                dt = datetime.strptime(clean, fmt)
                return dt.replace(tzinfo=timezone.utc).timestamp()
            except:
                continue
    except:
        pass
    return None

def _classify_post(post):
    t = str(_get_first(post, ["type","productType","product_type"], "")).lower()
    if "video" in t or "reel" in t or "clip" in t:
        return "reels"
    elif "carousel" in t or "sidecar" in t or "album" in t:
        return "carousels"
    # If from reel scraper, always classify as reel
    if post.get("_source") == "reel_scraper":
        return "reels"
    return "images"

def compute_day_distribution(parsed):
    """
    Works out which day of the week the brand posts most/least on,
    based on each post's real UTC date (not Instagram's pinned-post
    display order, and not local server time).

    Returns most_active_day and least_active_day. If a weekday never
    appears in the sample, it's reported as "never posted" rather than
    just being called the least-active day, since that's more accurate
    than implying it *did* happen, just rarely.
    """
    day_names = ["Monday","Tuesday","Wednesday","Thursday",
                 "Friday","Saturday","Sunday"]
    counts = Counter()
    for p in parsed:
        if p["ts"]:
            weekday = datetime.fromtimestamp(p["ts"], tz=timezone.utc).strftime("%A")
            counts[weekday] += 1

    if not counts:
        return {"most_active_day": "N/A", "least_active_day": "N/A",
                "best_performing_day": "N/A"}

    most_active_day = max(counts, key=counts.get)
    zero_days = [d for d in day_names if d not in counts]
    if zero_days:
        least_active_day = f"{zero_days[0]} (never posted)"
    else:
        least_active_day = min(counts, key=counts.get)

    # ── Best performing day: highest AVERAGE engagement per weekday.
    #    Requires 2+ posts on that weekday to qualify (same anti-fluke
    #    rule as best_performing_hour); falls back to the single best
    #    day if nothing qualifies. ─────────────────────────────────────
    day_posts = {}
    for p in parsed:
        if p["ts"]:
            wd = datetime.fromtimestamp(p["ts"], tz=timezone.utc).strftime("%A")
            day_posts.setdefault(wd, []).append(p)
    qualifying = {d: ps for d, ps in day_posts.items() if len(ps) >= 2}
    pool = qualifying if qualifying else day_posts
    if pool:
        best_performing_day = max(pool, key=lambda d: sum(_score(pp) for pp in pool[d]) / len(pool[d]))
    else:
        best_performing_day = "N/A"

    return {"most_active_day": most_active_day,
            "least_active_day": least_active_day,
            "best_performing_day": best_performing_day}

def _score(p):
    return p["likes"] + p["comments"] + (p["shares"] or 0) + (p["reposts"] or 0)


def compute_hour_insights(parsed, min_posts_to_qualify=2, outlier_multiplier=3):
    """
    Most Frequent Hour: which hour of day (UTC) has the most posts — pure count,
    doesn't consider engagement at all.

    Best Performing Hour: which hour has the highest AVERAGE engagement. Only
    hours with `min_posts_to_qualify` or more posts are allowed to compete —
    this stops a single lucky post from falsely looking like a repeatable
    "best time to post."

    Outlier note: separately, checks every hour that has exactly ONE post
    (too few to qualify above) and flags it if that single post scored
    dramatically higher than the sample's overall average (>= outlier_multiplier
    times the average). This keeps the standout post visible to the client
    without mislabeling it as a proven pattern.
    """
    hour_posts = {}
    for p in parsed:
        if p["ts"] is not None:
            hour = datetime.fromtimestamp(p["ts"], tz=timezone.utc).hour
            hour_posts.setdefault(hour, []).append(p)

    if not hour_posts:
        return {"most_frequent_hour": "N/A", "best_performing_hour": "N/A", "outlier_note": None}

    def hour_avg(posts):
        return sum(_score(pp) for pp in posts) / len(posts)

    # ── Most Frequent Hour — pure count ──────────────────────────────
    most_frequent_hour = max(hour_posts, key=lambda h: len(hour_posts[h]))

    # ── Best Performing Hour — avg engagement, min posts to qualify ──
    qualifying = {h: posts for h, posts in hour_posts.items() if len(posts) >= min_posts_to_qualify}
    if qualifying:
        best_hour = max(qualifying, key=lambda h: hour_avg(qualifying[h]))
        best_performing_hour = f"{best_hour:02d}:00 UTC"
    else:
        best_performing_hour = "N/A (not enough data)"

    # ── Outlier note — single-post hours only ────────────────────────
    all_scores  = [_score(pp) for pp in parsed]
    overall_avg = sum(all_scores) / len(all_scores) if all_scores else 0
    outlier_note = None
    lonely_hours = {h: posts[0] for h, posts in hour_posts.items() if len(posts) == 1}
    if lonely_hours and overall_avg > 0:
        lonely_hour, lonely_post = max(lonely_hours.items(), key=lambda kv: _score(kv[1]))
        lonely_score = _score(lonely_post)
        if lonely_score >= overall_avg * outlier_multiplier:
            outlier_note = (f"Outlier: one post at {lonely_hour:02d}:00 UTC scored {lonely_score:,} "
                             f"vs. an average of {overall_avg:.0f} — worth reviewing, though it's "
                             f"just one post so far.")

    return {
        "most_frequent_hour":   f"{most_frequent_hour:02d}:00 UTC",
        "best_performing_hour": best_performing_hour,
        "outlier_note":         outlier_note,
    }


def compute_posting_consistency(parsed):
    """
    Looks at the real gap (in days) between every consecutive pair of posts
    and measures how much those gaps vary, using the coefficient of variation
    (std deviation ÷ average gap). A low value means posts land at fairly
    even intervals ("Consistent"); a high value means bursts of activity
    followed by long silences ("Irregular") — something a single flat
    "X per week" number hides completely.
    """
    timestamps = sorted(p["ts"] for p in parsed if p["ts"] is not None)
    if len(timestamps) < 3:
        return "N/A"

    gaps = [(timestamps[i+1] - timestamps[i]) / 86400 for i in range(len(timestamps)-1)]
    avg_gap = sum(gaps) / len(gaps)
    if avg_gap == 0:
        return "N/A"
    variance = sum((g - avg_gap) ** 2 for g in gaps) / len(gaps)
    coefficient_of_variation = (variance ** 0.5) / avg_gap

    return "Consistent" if coefficient_of_variation < 0.75 else "Irregular"


def compute_momentum(parsed):
    """
    Splits posts into two halves by real date — the most recent half vs.
    the older half — and compares average engagement between them. Returns
    a signed percentage plus a direction ("up"/"down") so the report can
    show a colored arrow: green ▲ if recent posts are outperforming older
    ones, red ▼ if engagement is slipping.
    """
    dated = sorted((p for p in parsed if p["ts"] is not None), key=lambda p: p["ts"], reverse=True)
    n = len(dated)
    if n < 10:
        return {"momentum_pct": "N/A", "momentum_direction": None}

    half = n // 2
    recent_half = dated[:half]
    older_half  = dated[half:half*2] if len(dated) >= half*2 else dated[half:]
    if not older_half:
        return {"momentum_pct": "N/A", "momentum_direction": None}

    recent_avg = sum(_score(p) for p in recent_half) / len(recent_half)
    older_avg  = sum(_score(p) for p in older_half) / len(older_half)
    if older_avg == 0:
        return {"momentum_pct": "N/A", "momentum_direction": None}

    pct_change = (recent_avg - older_avg) / older_avg * 100
    return {
        "momentum_pct":       f"{pct_change:+.0f}%",
        "momentum_direction": "up" if pct_change >= 0 else "down",
    }

def parse_num(val):
    """
    Module-level numeric parser for follower counts etc. — handles
    values like "233,374", "1.2M", "45K", "1000+".

    NOTE: create_ppt also defines its own local parse_num (identical
    logic); that local copy harmlessly shadows this one inside
    create_ppt. This module-level version exists because
    fetch_linkedin needs it too, and previously calling it there
    raised a silent NameError (the local one wasn't visible outside
    create_ppt), which wiped out all LinkedIn post-level metrics.
    """
    try:
        v = str(val).replace(",","").upper().replace("+","")
        if "M" in v: return float(v.replace("M","")) * 1_000_000
        if "K" in v: return float(v.replace("K","")) * 1_000
        return float(v)
    except:
        return 0


def _smart_round(value, max_decimals=2):
    """
    Rounds to a whole number normally (e.g. 4.2 -> "4"), but if that
    would hide a real, non-zero value behind a flat "0" (e.g. 0.375
    likes/post), keeps adding decimal places until the number is
    actually visible — up to max_decimals. A genuinely-zero value still
    correctly shows "0". This fixes averages that looked misleadingly
    flat on low-volume post groups (e.g. 8 videos with only 3 total
    likes between them) while every other stat on the same slide showed
    a proper decimal.
    """
    if value == 0:
        return "0"
    rounded = round(value)
    if rounded != 0:
        return str(rounded)
    for d in range(1, max_decimals + 1):
        r = round(value, d)
        if r != 0:
            return f"{r:.{d}f}"
    return f"{value:.{max_decimals}f}"


def _compute_group_metrics(posts_group, followers, is_reels=False):
    n = len(posts_group)
    if n == 0:
        return {"n": 0}

    total_likes    = sum(p["likes"] for p in posts_group)
    total_comments = sum(p["comments"] for p in posts_group)
    shares_known   = [p["shares"] for p in posts_group if p["shares"] is not None]
    reposts_known  = [p["reposts"] for p in posts_group if p["reposts"] is not None]
    total_shares   = sum(shares_known) if shares_known else 0
    total_reposts  = sum(reposts_known) if reposts_known else 0

    avg_engagement  = (total_likes + total_comments + total_shares + total_reposts) / n
    er_per_follower = (total_likes + total_comments) / n / followers * 100 if followers else 0
    # Estimated ER: measured rate + 10% uplift to account for shares and
    # saves that platforms don't expose publicly. Shown on format slides
    # and used by the analysis; the uplift itself is NOT explained on the
    # slide (per client feedback) — it lives only here in the code.
    er_estimated = er_per_follower * 1.10

    scored = sorted(posts_group, key=_score, reverse=True)
    top    = scored[0]

    # ── Worst pick fairness rule ─────────────────────────────────
    # A post must be at least 3 days old to be eligible for "Worst".
    # A reel posted today has had almost no time to collect likes/
    # comments/shares, so its low score reflects its age, not its
    # quality — crowning it "worst" would mislead the client. The Top
    # pick is NOT filtered: if a brand-new post is already the top
    # scorer, it earned that early. If every post in the group is
    # newer than 3 days (very active accounts), fall back to judging
    # all of them, since there's nothing older to compare against.
    WORST_MIN_AGE_SECONDS = 3 * 86400
    now_ts = time.time()
    eligible_for_worst = [p for p in posts_group
                          if p.get("ts") is None or (now_ts - p["ts"]) >= WORST_MIN_AGE_SECONDS]
    if not eligible_for_worst:
        eligible_for_worst = posts_group
    worst = min(eligible_for_worst, key=_score)

    result = {
        "n":               n,
        "avg_engagement":  round(avg_engagement, 1),
        "er_per_follower": f"{er_per_follower:.2f}%",
        "er_estimated":    f"{er_estimated:.2f}%",
        "avg_likes":       _smart_round(total_likes / n),
        "avg_comments":    _smart_round(total_comments / n),
        "avg_shares":      _smart_round(total_shares / n) if shares_known else "N/A",
        "avg_reposts":     _smart_round(total_reposts / n) if reposts_known else "N/A",
        "top_score":       _score(top),
        "top_url":         top["url"],
        "worst_score":     _score(worst),
        "worst_url":       worst["url"],
    }

    if is_reels:
        views_known = [p["views"] for p in posts_group if p["views"] is not None]
        total_views = sum(views_known) if views_known else 0
        result["avg_views"] = _smart_round(total_views / n) if views_known else "N/A"
        if views_known:
            max_v = max(posts_group, key=lambda p: p["views"] or 0)
            result["max_views"]     = max_v["views"]
            result["max_views_url"] = max_v["url"]
        else:
            result["max_views"]     = "N/A"
            result["max_views_url"] = "N/A"

        # ── Rates are per VIEW (not per follower) ──────────────────
        result["like_rate"]    = f"{(total_likes/total_views*100):.2f}%"    if total_views else "N/A"
        result["comment_rate"] = f"{(total_comments/total_views*100):.2f}%" if total_views else "N/A"
        result["share_rate"]   = f"{(total_shares/total_views*100):.2f}%"   if (total_views and shares_known) else "N/A"
        result["repost_rate"]  = f"{(total_reposts/total_views*100):.2f}%"  if (total_views and reposts_known) else "N/A"
        result["er_by_view"]   = f"{((total_likes+total_comments+total_shares+total_reposts)/total_views*100):.2f}%" if total_views else "N/A"


    return result


def fetch_apify_async(username, actor_id, input_json):
    """Generic async Apify runner: start run, poll until done, fetch dataset."""
    try:
        print(f"   🔄 Apify async: starting {actor_id} for @{username}...")
        start_r = requests.post(
            f"https://api.apify.com/v2/acts/{actor_id}/runs",
            params={"token": APIFY_API_KEY},
            json=input_json,
            timeout=30
        )
        if start_r.status_code not in (200, 201):
            print(f"   ⚠️ Apify async start failed: {start_r.status_code}")
            return []

        run_id = start_r.json().get("data", {}).get("id")
        if not run_id:
            print("   ⚠️ Apify async: no run ID returned")
            return []

        print(f"   ⏳ Apify run started: {run_id} — polling...")

        last_status_r = None
        for attempt in range(24):
            time.sleep(5)
            last_status_r = requests.get(
                f"https://api.apify.com/v2/acts/{actor_id}/runs/{run_id}",
                params={"token": APIFY_API_KEY},
                timeout=15
            )
            status = last_status_r.json().get("data", {}).get("status", "")
            print(f"   ⏳ Apify status: {status} (attempt {attempt+1})")
            if status == "SUCCEEDED":
                break
            elif status in ("FAILED", "ABORTED", "TIMED-OUT"):
                print(f"   ⚠️ Apify run ended with: {status}")
                return []

        if last_status_r is None:
            return []

        dataset_id = last_status_r.json().get("data", {}).get("defaultDatasetId")
        if not dataset_id:
            print("   ⚠️ Apify async: no dataset ID")
            return []

        items_r = requests.get(
            f"https://api.apify.com/v2/datasets/{dataset_id}/items",
            params={"token": APIFY_API_KEY, "limit": 35},
            timeout=30
        )
        posts = items_r.json()
        if isinstance(posts, list) and len(posts) > 0:
            print(f"   ✅ Apify async returned {len(posts)} items")
            return posts
        else:
            print("   ⚠️ Apify async dataset was empty")
            return []

    except Exception as e:
        print(f"   ⚠️ Apify async error: {e}")
        return []


def fetch_post_metrics_via_analytics_actor(post_urls):
    """
    Enriches posts with shares / reposts / saves using the Apify actor
    'patient_discovery/instagram-reel-analytics-by-url'.

    Why this specific actor: after testing 6 different tools (SearchAPI,
    two official Apify scrapers, ScraperAPI, SocialCrawl, anonymous
    GraphQL), this was the ONLY one that returned real values. Verified
    against a live TalkingLands post on 2026-07-15:
        metrics.repost_count = 17   (matched browser exactly)
        metrics.share_count  = 750
        metrics.save_count   = 410
    Instagram only exposes these numbers to logged-in viewers; this actor
    handles that on its side, so no Instagram credentials are needed here.

    Takes a list of post URLs (supports bulk input), returns a dict of
    {url: {"shares": int|None, "reposts": int|None, "saves": int|None}}.
    Cost: one Apify result per post (~$2.30-2.70 per 1,000 → roughly
    ₹6-7 extra per 30-post report).
    """
    if not APIFY_API_KEY:
        print("   ⚠️ No APIFY_API_KEY set — skipping shares/reposts enrichment")
        return {}

    urls = [u for u in post_urls if u and u != "N/A"]
    if not urls:
        return {}

    results = {}
    try:
        print(f"   🔄 Fetching shares/reposts/saves for {len(urls)} posts via analytics actor...")
        actor_input = {"urls": urls}
        r = requests.post(
            "https://api.apify.com/v2/acts/patient_discovery~instagram-reel-analytics-by-url/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 180},
            json=actor_input,
            timeout=200
        )
        items = []
        if r.status_code in (200, 201):
            items = r.json()
            if not isinstance(items, list):
                items = []
        if not items:
            print(f"   ⚠️ Analytics actor sync empty/failed (status {r.status_code}), trying async...")
            items = fetch_apify_async("bulk", "patient_discovery~instagram-reel-analytics-by-url", actor_input)

        for item in items:
            metrics = item.get("metrics", {}) or {}
            # Match result back to its post URL: the actor echoes the input
            # URL and also provides the shortcode in `code`.
            item_url = _get_first(item, ["inputUrl", "url"]) or ""
            code     = item.get("code", "")
            key = None
            for u in urls:
                if (item_url and item_url.rstrip("/") == u.rstrip("/")) or (code and f"/{code}" in u):
                    key = u
                    break
            if key is None:
                continue
            results[key] = {
                "shares":  metrics.get("share_count"),
                "reposts": metrics.get("repost_count"),
                "saves":   metrics.get("save_count"),
            }

        found = sum(1 for v in results.values() if v.get("reposts") is not None or v.get("shares") is not None)
        print(f"   ✅ Shares/reposts found for {found}/{len(urls)} posts")
    except Exception as e:
        print(f"   ⚠️ Analytics actor error: {e}")

    return results


def fetch_instagram_posts_via_apify(username):
    """
    Two-scraper approach:
    1. apify/instagram-post-scraper  → images + carousels (likes, comments, views)
    2. apify/instagram-reel-scraper  → reels (likes, comments, views, SHARES)

    NOTE on the 35-post fetch limit (up from 30):
    Instagram displays pinned posts first in profile grid order, regardless of
    when they were actually posted — a post pinned 3 months ago can occupy a
    "top 30" slot and silently push out a genuinely recent post. We fetch a
    small buffer of extra posts here (35 instead of 30) so that, after the
    real-date sort applied later in fetch_instagram(), there's always enough
    genuinely recent posts to fill a true top-30 even when a few pinned posts
    turn out to be old. We intentionally do NOT slice down to 30 here — the
    final trim happens after sorting by real post date, not fetch order.
    """
    if not APIFY_API_KEY:
        print("   ⚠️ No APIFY_API_KEY set")
        return []

    all_posts = []

    # ── 1. Post scraper (images + carousels + some reels) ────────
    try:
        print(f"   🔄 Apify posts: fetching for @{username}...")
        post_input = {
            "username":      [username],
            "resultsLimit":  35,
            "addParentData": False,
        }
        r = requests.post(
            "https://api.apify.com/v2/acts/apify~instagram-post-scraper/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 120},
            json=post_input,
            timeout=150
        )
        if r.status_code in (200, 201):
            posts = r.json()
            if isinstance(posts, list) and len(posts) > 0:
                print(f"   ✅ Post scraper returned {len(posts)} posts")
                all_posts.extend(posts)
            else:
                print("   ⚠️ Post scraper sync empty, trying async...")
                posts = fetch_apify_async(username, "apify~instagram-post-scraper", post_input)
                all_posts.extend(posts)
        else:
            print(f"   ⚠️ Post scraper status {r.status_code}, trying async...")
            posts = fetch_apify_async(username, "apify~instagram-post-scraper", post_input)
            all_posts.extend(posts)
    except Exception as e:
        print(f"   ⚠️ Post scraper error: {e}")

    # ── 2. Reel scraper (reels with shares) ──────────────────────
    try:
        print(f"   🔄 Apify reels: fetching reels with shares for @{username}...")
        reel_input = {
            "username":     [username],
            "resultsLimit": 35,
        }
        r2 = requests.post(
            "https://api.apify.com/v2/acts/apify~instagram-reel-scraper/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 120},
            json=reel_input,
            timeout=150
        )
        if r2.status_code in (200, 201):
            reels = r2.json()
            if isinstance(reels, list) and len(reels) > 0:
                print(f"   ✅ Reel scraper returned {len(reels)} reels")
                for reel in reels:
                    reel["_source"] = "reel_scraper"
                existing_urls = {p.get("url","") for p in all_posts}
                new_reels = [r for r in reels if r.get("url","") not in existing_urls]
                all_posts.extend(new_reels)
                print(f"   ✅ Added {len(new_reels)} new reels (deduped)")
            else:
                print("   ⚠️ Reel scraper sync empty, trying async...")
                reels = fetch_apify_async(username, "apify~instagram-reel-scraper", reel_input)
                for reel in reels:
                    reel["_source"] = "reel_scraper"
                existing_urls = {p.get("url","") for p in all_posts}
                new_reels = [r for r in reels if r.get("url","") not in existing_urls]
                all_posts.extend(new_reels)
        else:
            print(f"   ⚠️ Reel scraper status {r2.status_code}, trying async...")
            reels = fetch_apify_async(username, "apify~instagram-reel-scraper", reel_input)
            for reel in reels:
                reel["_source"] = "reel_scraper"
            existing_urls = {p.get("url","") for p in all_posts}
            new_reels = [r for r in reels if r.get("url","") not in existing_urls]
            all_posts.extend(new_reels)
    except Exception as e:
        print(f"   ⚠️ Reel scraper error: {e}")

    print(f"   ✅ Total combined posts fetched (pre date-sort): {len(all_posts)}")
    return all_posts


def fetch_facebook_posts_via_apify(page_handle):
    """
    Fetches recent Facebook posts using apify/facebook-posts-scraper (the
    official Apify actor for Facebook Pages).

    HONEST CAVEAT: unlike the Instagram actors in this file, which went
    through many rounds of real-world testing to confirm exact field
    names, this Facebook actor has NOT yet been verified against real
    output. The field-name guesses below are based on common Apify
    Facebook-scraper conventions, with a wide fallback list per field
    (same defensive pattern used for Instagram). Use the /debug_facebook
    route to run one real test and confirm/adjust the field names before
    trusting this at scale — same process that eventually nailed Instagram.
    """
    if not APIFY_API_KEY:
        print("   ⚠️ No APIFY_API_KEY set")
        return []

    page_url = page_handle if page_handle.startswith("http") else f"https://www.facebook.com/{page_handle}"
    fb_input = {
        "startUrls": [{"url": page_url}],
        "resultsLimit": 35,  # same buffer logic as Instagram — covers pinned-post displacement
    }
    try:
        print(f"   🔄 Apify Facebook posts: fetching for {page_handle}...")
        r = requests.post(
            "https://api.apify.com/v2/acts/apify~facebook-posts-scraper/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 120},
            json=fb_input,
            timeout=150
        )
        if r.status_code in (200, 201):
            posts = r.json()
            if isinstance(posts, list) and len(posts) > 0:
                print(f"   ✅ Facebook post scraper returned {len(posts)} posts")
                return posts
        print(f"   ⚠️ Facebook post scraper sync empty/failed (status {r.status_code}), trying async...")
        return fetch_apify_async(page_handle, "apify~facebook-posts-scraper", fb_input)
    except Exception as e:
        print(f"   ⚠️ Facebook post scraper error: {e}")
        return []


def _classify_facebook_post(post):
    """
    Facebook doesn't have a Carousel format like Instagram, so posts are
    split into Photos / Videos / Links instead.

    CONFIRMED against real apify/facebook-posts-scraper output (2026-07-16):
    there is no top-level "type" field on a post. The real signal lives
    inside the "media" array — each media item has a "__typename" field
    ("Photo" for images; expected "Video" for actual video/reel posts,
    unconfirmed against a real video post yet). This replaced an earlier
    guess that checked for a top-level type field, which doesn't exist
    and was silently classifying every post as "photos".
    """
    media_items = post.get("media") or []
    typenames = [str(m.get("__typename","")).lower() for m in media_items if isinstance(m, dict)]
    if any("video" in t for t in typenames):
        return "videos"
    if any("photo" in t or "image" in t for t in typenames):
        return "photos"
    # No media at all, or an unrecognized shape — check the caption as a
    # last resort (posts often say "reel" or "video" somewhere in their
    # own text, e.g. "In this reel, we break down..."), then fall back
    # to "links" for pure text/link posts.
    text = str(post.get("text","")).lower()
    if "reel" in text or "video" in text:
        return "videos"
    if not media_items:
        return "links"
    return "photos"


def fetch_instagram(username):
    print(f"📸 Fetching Instagram: @{username}")
    data = {}
    try:
        # ── Profile data from SearchAPI — WITH RETRIES, NEVER FATAL ───
        # Previously: a single failed profile call did `return data` and
        # abandoned Instagram entirely, producing a report where every
        # Instagram slide showed N/A even though the Apify post scrapers
        # were healthy (confirmed live 2026-07-18: FB/YT/LI all worked in
        # the same run). Now: retry up to 3 times with a short pause, and
        # if the profile STILL can't be fetched, continue with post-level
        # scraping anyway — followers show N/A but reels/images/carousels
        # data, top/worst picks, and view-based rates still populate.
        p = {}
        ig = {}
        for attempt in range(3):
            try:
                r = requests.get(
                    "https://www.searchapi.io/api/v1/search",
                    params={"engine":"instagram_profile","username":username,
                            "api_key":SEARCHAPI_KEY},
                    timeout=30
                )
                if r.status_code == 200:
                    ig = r.json()
                    p  = ig.get("profile", {}) or {}
                    if p:
                        break
                print(f"   ⚠️ Instagram profile attempt {attempt+1} failed "
                      f"(status {r.status_code}) — retrying..." if attempt < 2 else
                      f"   ⚠️ Instagram profile attempt {attempt+1} failed (status {r.status_code})")
            except Exception as e:
                print(f"   ⚠️ Instagram profile attempt {attempt+1} error: {e}")
            time.sleep(3)
        if not p:
            print("   ⚠️ Instagram profile unavailable after 3 attempts — "
                  "continuing with post-level data only")

        followers   = p.get("followers", "N/A")
        following   = p.get("following", "N/A")
        total_posts = p.get("posts", "N/A")
        bio         = p.get("biography", "")
        full_name   = p.get("full_name", "")
        is_verified = p.get("is_verified", False)
        is_business = p.get("is_business", False)
        avatar_url  = _get_first(p, ["avatar_hd", "avatar"])

        try:
            followers_n = int(str(followers).replace(",",""))
        except:
            followers_n = 0

        # ── Posts via Apify (post + reel scrapers), fallback to SearchAPI ──
        apify_posts = fetch_instagram_posts_via_apify(username)
        if apify_posts:
            posts_raw = apify_posts
            source    = "apify"
        else:
            posts_raw = ig.get("posts", [])
            source    = "searchapi"
            print(f"   ↩️  Using SearchAPI posts ({len(posts_raw)} posts)")

        print(f"   🔍 Source: {source} | Posts fetched: {len(posts_raw)}")

        # ── Parse every post ──────────────────────────────────────────
        parsed = []

        for post in posts_raw:
            if source == "apify":
                reel_src = post.get("_source") == "reel_scraper"
                likes    = post.get("likesCount", 0) or 0
                comments = post.get("commentsCount", 0) or 0
                if reel_src:
                    # Reel scraper has reliable sharesCount
                    shares  = _get_first(post, ["sharesCount","videoSharesCount","shares"])
                    reposts = _get_first(post, ["repostsCount","reposts"])
                    views   = _get_first(post, ["videoViewCount","videoPlayCount","playsCount","views"])
                else:
                    shares  = _get_first(post, ["sharesCount","share_count","shares"])
                    reposts = _get_first(post, ["repostsCount","repost_count","reposts"])
                    views   = _get_first(post, ["videoViewCount","viewsCount","videoPlayCount","views"])
                url = _get_first(post, ["url","permalink","link"]) or "N/A"
                ts  = _to_timestamp(_get_first(post, ["timestamp","iso_date"]))
            else:
                # SearchAPI confirmed field names
                likes    = post.get("likes", 0) or 0
                comments = post.get("comments", 0) or 0
                shares   = _get_first(post, ["share_count","shares","reshare_count"])
                reposts  = _get_first(post, ["repost_count","reposts","reshares"])
                views    = _get_first(post, ["views","view_count","play_count"])
                url      = _get_first(post, ["permalink","link","url","post_url"]) or "N/A"
                ts       = _to_timestamp(_get_first(post, ["iso_date","timestamp","taken_at","created_at"]))

            ptype = _classify_post(post)

            parsed.append({
                "type":    ptype,
                "likes":   likes,
                "comments":comments,
                "shares":  (int(shares)  if shares  is not None else None),
                "reposts": (int(reposts) if reposts is not None else None),
                "views":   (int(views)   if views   is not None else None),
                "url":     url,
                "ts":      ts,
            })

        # ── Sort by REAL post date, ignoring Instagram's pinned-post
        #    display order, and keep only the true most-recent 30.
        #
        #    Why: Instagram shows pinned posts first in the grid no matter
        #    how old they are. If we trusted fetch/display order, an old
        #    pinned post could occupy a "top 30" slot and quietly displace
        #    a genuinely recent post — skewing posting frequency, averages,
        #    and which post gets crowned Top/Worst performer. A pinned post
        #    that IS genuinely recent is unaffected by this — it earns its
        #    place in the top 30 purely by having a recent real date, same
        #    as any other post. Pin status itself is never checked. ───────
        posts_with_dates    = [pp for pp in parsed if pp["ts"] is not None]
        posts_without_dates = [pp for pp in parsed if pp["ts"] is None]
        posts_with_dates.sort(key=lambda pp: pp["ts"], reverse=True)
        parsed = (posts_with_dates + posts_without_dates)[:30]

        # ── Enrich with shares / reposts / saves via the analytics actor ──
        #    Runs on the final 30 posts only (after the date-sort/trim), so
        #    we never pay for posts that won't appear in the report. Values
        #    already present from the base scrapers are kept; the actor
        #    fills the gaps.
        posts_needing = [pp for pp in parsed
                         if pp["reposts"] is None or pp["shares"] is None]
        if posts_needing and APIFY_API_KEY:
            metric_results = fetch_post_metrics_via_analytics_actor(
                [pp["url"] for pp in posts_needing]
            )
            for pp in parsed:
                found = metric_results.get(pp["url"])
                if not found:
                    continue
                if pp["reposts"] is None and found.get("reposts") is not None:
                    pp["reposts"] = int(found["reposts"])
                if pp["shares"] is None and found.get("shares") is not None:
                    pp["shares"] = int(found["shares"])

        # ── Content-type counts — computed AFTER the date-sort/trim above,
        #    so the Content Strategy pie chart reflects the true final 30
        #    posts, not the larger raw fetch batch. ───────────────────────
        img_c = sum(1 for pp in parsed if pp["type"] == "images")
        car_c = sum(1 for pp in parsed if pp["type"] == "carousels")
        vid_c = sum(1 for pp in parsed if pp["type"] == "reels")

        post_count      = len(parsed)
        images_group    = [pp for pp in parsed if pp["type"] == "images"]
        carousels_group = [pp for pp in parsed if pp["type"] == "carousels"]
        reels_group     = [pp for pp in parsed if pp["type"] == "reels"]

        # ── Overall engagement ────────────────────────────────────────
        total_l          = sum(pp["likes"] for pp in parsed)
        total_c          = sum(pp["comments"] for pp in parsed)
        engagement_total = total_l + total_c
        eng_rate         = f"{(engagement_total/post_count/followers_n*100):.2f}%" if (post_count and followers_n) else "N/A"

        reel_l = sum(pp["likes"]    for pp in reels_group)
        reel_c = sum(pp["comments"] for pp in reels_group)
        eng_rate_reels = (f"{(reel_l+reel_c)/len(reels_group)/followers_n*100:.2f}%"
                          if (reels_group and followers_n) else "N/A")

        # ── Posting frequency ─────────────────────────────────────────
        timestamps = [pp["ts"] for pp in parsed if pp["ts"]]
        if len(timestamps) >= 2:
            span_days         = max((max(timestamps) - min(timestamps)) / 86400, 1)
            posting_frequency = f"{(post_count / span_days * 7):.1f} / week"
        else:
            posting_frequency = "N/A"

        # ── Day-of-week posting distribution (UTC) ──────────────────────
        day_dist = compute_day_distribution(parsed)

        # ── Best/most-frequent posting hour, consistency, momentum (UTC) ──
        hour_insights = compute_hour_insights(parsed)
        consistency   = compute_posting_consistency(parsed)
        momentum      = compute_momentum(parsed)

        # ── Per-group metrics ─────────────────────────────────────────
        images_metrics    = _compute_group_metrics(images_group,    followers_n)
        carousels_metrics = _compute_group_metrics(carousels_group, followers_n)
        reels_metrics     = _compute_group_metrics(reels_group,     followers_n, is_reels=True)

        data = {
            "username":    username,
            "full_name":   full_name,
            "followers":   str(followers),
            "following":   str(following),
            "posts":       str(total_posts),
            "bio":         bio,
            "is_verified": "Yes" if is_verified else "No",
            "is_business": "Yes" if is_business else "No",

            "sample_size":           post_count,
            "engagement_rate":       eng_rate,
            "engagement_rate_reels": eng_rate_reels,
            "engagement_total":      engagement_total,
            "posting_frequency":     posting_frequency,

            "most_active_day":  day_dist["most_active_day"],
            "least_active_day": day_dist["least_active_day"],
            "best_performing_day": day_dist["best_performing_day"],
            "avg_engagement": round(engagement_total / post_count, 1) if post_count else "N/A",

            "avatar_url": avatar_url,

            "most_frequent_hour":   hour_insights["most_frequent_hour"],
            "best_performing_hour": hour_insights["best_performing_hour"],
            "outlier_note":         hour_insights["outlier_note"],
            "posting_consistency":  consistency,
            "momentum_pct":         momentum["momentum_pct"],
            "momentum_direction":   momentum["momentum_direction"],

            "img_count":  img_c,
            "car_count":  car_c,
            "vid_count":  vid_c,
            "post_count": post_count,

            "images":    images_metrics,
            "carousels": carousels_metrics,
            "reels":     reels_metrics,
        }
        print(f"   ✅ {followers} followers | {post_count} posts "
              f"({img_c} img / {car_c} carousel / {vid_c} reels) | ER: {eng_rate} | "
              f"Most active: {day_dist['most_active_day']} | Least active: {day_dist['least_active_day']}")
    except Exception as e:
        import traceback
        print(f"   ⚠️ Instagram failed: {e}")
        print(traceback.format_exc())
    return data


def fetch_facebook(username):
    print(f"👥 Fetching Facebook: {username}")
    data = {}
    try:
        # ── Page data from SearchAPI — with retries, never fatal (same
        #    hardening as Instagram: a failed profile call must not wipe
        #    out the whole Facebook section when Apify posts still work) ──
        p = {}
        for attempt in range(3):
            try:
                r = requests.get(
                    "https://www.searchapi.io/api/v1/search",
                    params={"engine":"facebook_business_page","username":username,
                            "api_key":SEARCHAPI_KEY},
                    timeout=30
                )
                if r.status_code == 200:
                    p = r.json().get("page", {}) or {}
                    if p:
                        break
                print(f"   ⚠️ Facebook page attempt {attempt+1} failed (status {r.status_code})")
            except Exception as e:
                print(f"   ⚠️ Facebook page attempt {attempt+1} error: {e}")
            time.sleep(3)
        if not p:
            print("   ⚠️ Facebook page unavailable after 3 attempts — "
                  "continuing with post-level data only")

        followers_raw = p.get("followers",{}).get("count","N/A")
        try:
            followers_n = int(str(followers_raw).replace(",",""))
        except:
            followers_n = 0

        data = {
            "name":        p.get("name", username),
            "followers":   str(followers_raw),
            "following":   str(p.get("following",{}).get("count","N/A")),
            "category":    ", ".join(p.get("about",{}).get("category_formatted",[]) or []) or "N/A",
            "about":       p.get("about",{}).get("description","") or p.get("about",{}).get("general_info",""),
            "address":     p.get("address","") or "",
            "phone":       p.get("phone","") or "",
            "email":       p.get("email","") or "",
            "website":     p.get("website","") or "",
            "rating":      str(p.get("ratings",{}).get("value","N/A")),
            "is_verified": "Yes" if p.get("is_verified") else "No",
        }
        print(f"   ✅ {data['followers']} followers")

        # ── Posts via Apify — same pattern as Instagram ─────────────────
        posts_raw = fetch_facebook_posts_via_apify(username)
        print(f"   🔍 Facebook posts fetched: {len(posts_raw)}")

        parsed = []
        for post in posts_raw:
            likes    = _get_first(post, ["likes", "likesCount", "reactionsCount", "reactions"], 0) or 0
            comments = _get_first(post, ["comments", "commentsCount"], 0) or 0
            shares   = _get_first(post, ["shares", "sharesCount"])
            views    = _get_first(post, ["viewsCount", "views", "videoViewCount"])
            url      = _get_first(post, ["url", "postUrl", "facebookUrl", "link"]) or "N/A"
            ts       = _to_timestamp(_get_first(post, ["time", "timestamp", "date", "publishedAt"]))
            ptype    = _classify_facebook_post(post)
            try: likes = int(likes)
            except: likes = 0
            try: comments = int(comments)
            except: comments = 0
            parsed.append({
                "type":     ptype,
                "likes":    likes,
                "comments": comments,
                "shares":   (int(shares) if shares is not None else None),
                "reposts":  None,  # Facebook has no repost concept — always None, reused
                                    # formulas treat this as 0 automatically, same as Instagram's N/A handling
                "views":    (int(views) if views is not None else None),
                "url":      url,
                "ts":       ts,
            })

        # ── Sort by REAL post date, keep the true most-recent 30 —
        #    identical logic to the Instagram pinned-post fix. ──────────
        posts_with_dates    = [pp for pp in parsed if pp["ts"] is not None]
        posts_without_dates = [pp for pp in parsed if pp["ts"] is None]
        posts_with_dates.sort(key=lambda pp: pp["ts"], reverse=True)
        parsed = (posts_with_dates + posts_without_dates)[:30]

        photos_group = [pp for pp in parsed if pp["type"] == "photos"]
        videos_group = [pp for pp in parsed if pp["type"] == "videos"]
        links_group  = [pp for pp in parsed if pp["type"] == "links"]
        post_count   = len(parsed)

        total_l = sum(pp["likes"] for pp in parsed)
        total_c = sum(pp["comments"] for pp in parsed)
        engagement_total = total_l + total_c
        eng_rate = f"{(engagement_total/post_count/followers_n*100):.2f}%" if (post_count and followers_n) else "N/A"

        vid_l = sum(pp["likes"] for pp in videos_group)
        vid_c = sum(pp["comments"] for pp in videos_group)
        eng_rate_videos = (f"{(vid_l+vid_c)/len(videos_group)/followers_n*100:.2f}%"
                            if (videos_group and followers_n) else "N/A")

        timestamps = [pp["ts"] for pp in parsed if pp["ts"]]
        if len(timestamps) >= 2:
            span_days = max((max(timestamps) - min(timestamps)) / 86400, 1)
            posting_frequency = f"{(post_count / span_days * 7):.1f} / week"
        else:
            posting_frequency = "N/A"

        # ── Reused directly from the Instagram pipeline — same formulas,
        #    just fed Facebook's parsed posts instead. ───────────────────
        day_dist      = compute_day_distribution(parsed)
        hour_insights = compute_hour_insights(parsed)
        consistency   = compute_posting_consistency(parsed)
        momentum      = compute_momentum(parsed)

        photos_metrics = _compute_group_metrics(photos_group, followers_n)
        videos_metrics = _compute_group_metrics(videos_group, followers_n, is_reels=True)
        links_metrics  = _compute_group_metrics(links_group,  followers_n)

        data.update({
            "sample_size":            post_count,
            "engagement_rate":        eng_rate,
            "engagement_rate_videos": eng_rate_videos,
            "engagement_total":       engagement_total,
            "posting_frequency":      posting_frequency,

            "most_active_day":  day_dist["most_active_day"],
            "least_active_day": day_dist["least_active_day"],
            "best_performing_day": day_dist["best_performing_day"],
            "avg_engagement": round(engagement_total / post_count, 1) if post_count else "N/A",

            "most_frequent_hour":   hour_insights["most_frequent_hour"],
            "best_performing_hour": hour_insights["best_performing_hour"],
            "outlier_note":         hour_insights["outlier_note"],
            "posting_consistency":  consistency,
            "momentum_pct":         momentum["momentum_pct"],
            "momentum_direction":   momentum["momentum_direction"],

            "photo_count": len(photos_group),
            "video_count": len(videos_group),
            "link_count":  len(links_group),

            "photos": photos_metrics,
            "videos": videos_metrics,
            "links":  links_metrics,
        })
        print(f"   ✅ {post_count} posts ({len(photos_group)} photos / {len(videos_group)} videos / "
              f"{len(links_group)} links) | ER: {eng_rate}")

    except Exception as e:
        import traceback
        print(f"   ⚠️ Facebook failed: {e}")
        print(traceback.format_exc())
    return data


def _parse_iso8601_duration(duration_str):
    """
    Parses YouTube's ISO 8601 duration format (e.g. "PT4M13S", "PT58S",
    "PT1H2M3S") into total seconds. This is a real, stable, documented
    format from Google's own API — no guessing required here, unlike
    the scraper-based platforms.
    """
    if not duration_str:
        return 0
    match = re.match(r'PT(?:(\d+)H)?(?:(\d+)M)?(?:(\d+)S)?', duration_str)
    if not match:
        return 0
    hours   = int(match.group(1) or 0)
    minutes = int(match.group(2) or 0)
    seconds = int(match.group(3) or 0)
    return hours * 3600 + minutes * 60 + seconds


def _classify_youtube_video(video):
    """
    FALLBACK ONLY — used when the live /shorts/ redirect check (below)
    fails for a specific video, e.g. a network hiccup. This is NOT the
    primary classification method anymore: duration alone is wrong,
    because YouTube now allows Shorts up to 3 minutes long, but plenty
    of normal landscape videos are also under 3 minutes and are NOT
    Shorts (confirmed live: two TalkingLands videos at 2:24 and 2:22
    are regular landscape uploads, not Shorts). Real classification
    lives in _classify_youtube_videos_batch below.
    """
    duration_str = video.get("contentDetails", {}).get("duration", "")
    seconds = _parse_iso8601_duration(duration_str)
    return "shorts" if seconds <= 180 else "videos"


def _check_is_short_via_redirect(video_id, timeout=6):
    """
    Asks YouTube itself whether a video is a Short — this is the real
    fix for the duration-based misclassification bug.

    How it works: every Short also has a URL at
    https://www.youtube.com/shorts/{video_id}.
      - YouTube serves that URL directly (status 200)  -> it IS a Short
      - YouTube redirects it to /watch?v=... (30x)      -> it's a
        regular Video, even if it happens to be short in length

    This checks YouTube's own classification instead of guessing from
    duration, which is what was misclassifying short landscape videos
    (e.g. 2:22, 2:24 long) as Shorts.

    Returns True (Short), False (Video), or None if the request itself
    failed (network issue) — the caller falls back to duration only in
    that None case, never as the default.
    """
    url = f"https://www.youtube.com/shorts/{video_id}"
    headers = {"User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) Chrome/120.0.0.0"}
    try:
        r = requests.head(url, allow_redirects=False, timeout=timeout, headers=headers)
        # A few edge servers respond oddly to HEAD — retry with GET
        # once before giving up on this video.
        if r.status_code not in (200, 301, 302, 303, 307, 308):
            r = requests.get(url, allow_redirects=False, timeout=timeout, headers=headers)
        if r.status_code == 200:
            return True
        if r.status_code in (301, 302, 303, 307, 308):
            return False
    except Exception as e:
        print(f"   ⚠️ Shorts-check failed for {video_id}: {e}")
    return None


def _classify_youtube_videos_batch(videos, max_workers=10):
    """
    Classifies every video as Shorts vs Videos by checking YouTube's
    real /shorts/{id} redirect behavior for each one, run CONCURRENTLY
    so ~30 videos don't add ~30 sequential network round-trips to
    report generation time.

    Falls back to the old duration-based rule ONLY for a video whose
    live check genuinely failed (e.g. a network hiccup) — never used
    as the default/primary method anymore.
    """
    video_ids = [v.get("id","") for v in videos if v.get("id")]
    redirect_results = {}
    if video_ids:
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_id = {executor.submit(_check_is_short_via_redirect, vid): vid
                             for vid in video_ids}
            for future in concurrent.futures.as_completed(future_to_id):
                vid = future_to_id[future]
                try:
                    redirect_results[vid] = future.result()
                except Exception:
                    redirect_results[vid] = None

    classifications = {}
    fallback_used = 0
    for v in videos:
        vid = v.get("id", "")
        is_short = redirect_results.get(vid)
        if is_short is None:
            classifications[vid] = _classify_youtube_video(v)
            fallback_used += 1
        else:
            classifications[vid] = "shorts" if is_short else "videos"
    if fallback_used:
        print(f"   ⚠️ Shorts/Video: fell back to duration for {fallback_used} video(s) (live check failed)")
    return classifications


def fetch_youtube_videos_via_api(channel_handle_or_id):
    """
    Official YouTube Data API v3 — three-step fetch, all confirmed from
    Google's own published documentation, not a reverse-engineered
    scraper:
      1. channels.list  -> resolve handle/ID to the channel's uploads
                            playlist ID (costs 1 quota unit)
      2. playlistItems.list -> list up to 35 recent video IDs from that
                                playlist (costs 1 unit)
      3. videos.list    -> batch-fetch full stats for those videos in
                            ONE call (costs 1 unit, up to 50 IDs)
    Total: ~3 units per report, against a 10,000/day free quota.
    """
    if not YOUTUBE_API_KEY:
        print("   ⚠️ No YOUTUBE_API_KEY set")
        return [], None

    raw = channel_handle_or_id.strip().lstrip("@")

    try:
        if raw.startswith("UC") and len(raw) >= 20:
            ch_params = {"part": "snippet,contentDetails", "id": raw, "key": YOUTUBE_API_KEY}
        else:
            ch_params = {"part": "snippet,contentDetails", "forHandle": raw, "key": YOUTUBE_API_KEY}
        r = requests.get("https://www.googleapis.com/youtube/v3/channels", params=ch_params, timeout=30)
        items = r.json().get("items", [])
        if not items:
            print(f"   ⚠️ YouTube channel not found for: {channel_handle_or_id}")
            return [], None
        uploads_playlist_id = items[0]["contentDetails"]["relatedPlaylists"]["uploads"]
        # Prefer the real, readable handle (e.g. "@kekahr") over a raw
        # channel ID — some websites only expose the raw /channel/UC...
        # URL, which otherwise ends up displayed as-is on the slide.
        resolved_handle = items[0].get("snippet", {}).get("customUrl") or channel_handle_or_id
    except Exception as e:
        print(f"   ⚠️ YouTube channel resolution failed: {e}")
        return [], None

    try:
        r = requests.get("https://www.googleapis.com/youtube/v3/playlistItems", params={
            "part": "contentDetails", "playlistId": uploads_playlist_id,
            "maxResults": 35, "key": YOUTUBE_API_KEY
        }, timeout=30)
        video_ids = [item["contentDetails"]["videoId"] for item in r.json().get("items", [])]
        if not video_ids:
            print("   ⚠️ YouTube uploads playlist returned no videos")
            return [], resolved_handle
    except Exception as e:
        print(f"   ⚠️ YouTube playlist fetch failed: {e}")
        return [], resolved_handle

    try:
        r = requests.get("https://www.googleapis.com/youtube/v3/videos", params={
            "part": "snippet,statistics,contentDetails",
            "id": ",".join(video_ids), "key": YOUTUBE_API_KEY
        }, timeout=30)
        videos = r.json().get("items", [])
        print(f"   ✅ YouTube official API returned {len(videos)} videos")
        return videos, resolved_handle
    except Exception as e:
        print(f"   ⚠️ YouTube video stats fetch failed: {e}")
        return [], resolved_handle


def fetch_youtube(channel_id):
    print(f"▶️  Fetching YouTube: {channel_id}")
    data = {}
    raw = channel_id.strip().lstrip("@")
    if raw.startswith("UC") and len(raw) >= 20:
        param = raw
    else:
        param = "@" + raw
    print(f"   Using param: {param}")
    try:
        r = requests.get(
            "https://www.searchapi.io/api/v1/search",
            params={"engine":"youtube_channel","channel_id":param,
                    "api_key":SEARCHAPI_KEY}
        )
        if r.status_code == 200:
            result  = r.json()
            channel = result.get("channel", {})
            about   = result.get("about", {})
            data = {
                "handle":      channel.get("handle", channel_id),
                "title":       channel.get("title",""),
                "subscribers": str(channel.get("subscribers","N/A")),
                "videos":      str(channel.get("videos","N/A")),
                "views":       str(about.get("views","N/A")),
                "description": channel.get("description","") or about.get("description",""),
                "joined":      about.get("joined_date","N/A"),
                "is_verified": "Yes" if channel.get("is_verified") else "No",
            }
            print(f"   ✅ {data['subscribers']} subscribers | {data['videos']} videos")
    except Exception as e:
        print(f"   ⚠️ YouTube failed: {e}")

    # ── Videos via the OFFICIAL YouTube Data API v3 — reuses the exact
    #    same metric machinery as Instagram/Facebook/LinkedIn. ────────
    try:
        subscribers_n = parse_num(data.get("subscribers", 0))
        videos_raw, resolved_handle = fetch_youtube_videos_via_api(channel_id)
        if resolved_handle:
            # Overwrite with the real, readable handle if the API found
            # one — fixes cases where the website only exposed a raw
            # /channel/UC... URL instead of a friendly @handle.
            data["handle"] = resolved_handle
        print(f"   🔍 YouTube videos fetched: {len(videos_raw)}")

        # ── Classify ALL videos as Shorts vs. Videos in one concurrent
        #    batch, using YouTube's real /shorts/ redirect check instead
        #    of guessing from duration. Done once here, before the loop,
        #    so the loop below is just a plain dict lookup per video. ──
        yt_classifications = _classify_youtube_videos_batch(videos_raw)

        parsed = []
        for v in videos_raw:
            snippet = v.get("snippet", {}) or {}
            stats   = v.get("statistics", {}) or {}
            try: likes = int(stats.get("likeCount", 0) or 0)
            except: likes = 0
            try: comments = int(stats.get("commentCount", 0) or 0)
            except: comments = 0
            try: views = int(stats.get("viewCount", 0) or 0)
            except: views = 0
            video_id = v.get("id", "")
            url = f"https://www.youtube.com/watch?v={video_id}" if video_id else "N/A"
            ts = _to_timestamp(snippet.get("publishedAt"))
            ptype = yt_classifications.get(video_id, "videos")
            # No "shares" field exists anywhere in YouTube's official
            # API — genuinely doesn't exist, not a data gap on our end.
            parsed.append({
                "type": ptype, "likes": likes, "comments": comments,
                "shares": None, "reposts": None, "views": views,
                "url": url, "ts": ts,
            })

        posts_with_dates    = [pp for pp in parsed if pp["ts"] is not None]
        posts_without_dates = [pp for pp in parsed if pp["ts"] is None]
        posts_with_dates.sort(key=lambda pp: pp["ts"], reverse=True)
        parsed = (posts_with_dates + posts_without_dates)[:30]

        video_count = len(parsed)
        shorts_group = [pp for pp in parsed if pp["type"] == "shorts"]
        videos_group = [pp for pp in parsed if pp["type"] == "videos"]

        total_l = sum(pp["likes"] for pp in parsed)
        total_c = sum(pp["comments"] for pp in parsed)
        engagement_total = total_l + total_c
        eng_rate = f"{(engagement_total/video_count/subscribers_n*100):.2f}%" if (video_count and subscribers_n) else "N/A"

        timestamps = [pp["ts"] for pp in parsed if pp["ts"]]
        if len(timestamps) >= 2:
            span_days = max((max(timestamps) - min(timestamps)) / 86400, 1)
            posting_frequency = f"{(video_count / span_days * 7):.1f} / week"
        else:
            posting_frequency = "N/A"

        day_dist      = compute_day_distribution(parsed)
        hour_insights = compute_hour_insights(parsed)
        consistency   = compute_posting_consistency(parsed)
        momentum      = compute_momentum(parsed)

        # is_reels=True for BOTH groups — unlike Instagram/Facebook,
        # every YouTube video (Short or long-form) has a real view
        # count, so both categories get the view-based rate metrics.
        shorts_metrics = _compute_group_metrics(shorts_group, subscribers_n, is_reels=True)
        videos_metrics = _compute_group_metrics(videos_group, subscribers_n, is_reels=True)

        data.update({
            "sample_size":       video_count,
            "engagement_rate":   eng_rate,
            "engagement_total":  engagement_total,
            "posting_frequency": posting_frequency,
            "most_active_day":  day_dist["most_active_day"],
            "least_active_day": day_dist["least_active_day"],
            "best_performing_day": day_dist["best_performing_day"],
            "avg_engagement": round(engagement_total / video_count, 1) if video_count else "N/A",
            "avg_views_sample": (f"{sum((pp['views'] or 0) for pp in parsed) // video_count:,}"
                                  if video_count else "N/A"),
            "most_frequent_hour":   hour_insights["most_frequent_hour"],
            "best_performing_hour": hour_insights["best_performing_hour"],
            "outlier_note":         hour_insights["outlier_note"],
            "posting_consistency":  consistency,
            "momentum_pct":         momentum["momentum_pct"],
            "momentum_direction":   momentum["momentum_direction"],
            "shorts_count": len(shorts_group),
            "videos_count": len(videos_group),
            "shorts": shorts_metrics,
            "videos_performance": videos_metrics,
        })
        print(f"   ✅ {video_count} videos ({len(shorts_group)} shorts / {len(videos_group)} long-form) | ER: {eng_rate}")
    except Exception as e:
        import traceback
        print(f"   ⚠️ YouTube post-level fetch failed: {e}")
        print(traceback.format_exc())

    return data


def fetch_linkedin_posts_via_apify(company_handle):
    """
    Fetches recent LinkedIn company posts. Primary actor:
    brilliant_gum/linkedin-company-post-scraper — chosen because its
    documentation explicitly supports sorting by "recent" (chronological,
    newest first) and up to 500 posts per run.

    Fallback: data-slayer/linkedin-company-posts-scraper — already
    verified working against real data (field names confirmed), but it
    accepts ONLY a URL: it ignores any limit parameter (docs: "One
    field. Paste your company URLs and run."), which is why it
    returned just 10 posts. Kept as a safety net in case the primary
    actor's input guesses need adjusting.

    HONEST CAVEAT: brilliant_gum's exact input field names are
    best-effort guesses from its documentation (sort "recent"/"top",
    a post limit, pagination via start page) — not yet verified against
    a real run. Use /debug_linkedin to confirm. If the primary returns
    nothing, the proven fallback keeps reports working meanwhile.
    """
    if not APIFY_API_KEY:
        print("   ⚠️ No APIFY_API_KEY set")
        return []

    company_url = company_handle if company_handle.startswith("http") else f"https://www.linkedin.com/company/{company_handle}"

    # ── Primary: harvestapi/linkedin-company-posts ──────────────────
    # Input format CONFIRMED verbatim from Apify's official API example
    # for this actor: {"targetUrls": [...], "maxPosts": N}.
    # HarvestAPI: 5.8K users, 5.0 rating, $2/1k posts. Scrapes newest-
    # first by design ("from now up to postedLimitDate").
    # (Previous primary, brilliant_gum, was a 15-user unrated actor that
    # rejected every API run — likely needed manual rental activation.)
    primary_input = {
        "targetUrls": [company_url],
        "maxPosts":   35,
    }
    try:
        print(f"   🔄 Apify LinkedIn posts (harvestapi, newest-first): {company_handle}...")
        r = requests.post(
            "https://api.apify.com/v2/acts/harvestapi~linkedin-company-posts/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 120},
            json=primary_input,
            timeout=150
        )
        if r.status_code in (200, 201):
            posts = r.json()
            if isinstance(posts, list) and len(posts) > 0:
                print(f"   ✅ Primary LinkedIn scraper returned {len(posts)} posts")
                return posts
        print(f"   ⚠️ Primary LinkedIn scraper empty/failed (status {r.status_code}), falling back...")
    except Exception as e:
        print(f"   ⚠️ Primary LinkedIn scraper error: {e}, falling back...")

    # ── Fallback: data-slayer (proven field names, URL-only input) ──
    fallback_input = {"startUrls": [{"url": company_url}]}
    try:
        print(f"   🔄 Apify LinkedIn posts (fallback): {company_handle}...")
        r = requests.post(
            "https://api.apify.com/v2/acts/data-slayer~linkedin-company-posts-scraper/run-sync-get-dataset-items",
            params={"token": APIFY_API_KEY, "timeout": 120},
            json=fallback_input,
            timeout=150
        )
        if r.status_code in (200, 201):
            posts = r.json()
            if isinstance(posts, list) and len(posts) > 0:
                print(f"   ✅ Fallback LinkedIn scraper returned {len(posts)} posts")
                return posts
        print(f"   ⚠️ Fallback LinkedIn scraper sync empty/failed (status {r.status_code}), trying async...")
        return fetch_apify_async(company_handle, "data-slayer~linkedin-company-posts-scraper", fallback_input)
    except Exception as e:
        print(f"   ⚠️ Fallback LinkedIn scraper error: {e}")
        return []


def _classify_linkedin_post(post):
    """
    7 categories (Reposted Content deliberately excluded — a reposted
    video is just classified as a video, based on its actual attached
    format): Videos / Multi-Image / Single-Image / Documents / Polls /
    Text-Only / Other (Articles, Newsletters, Celebrations combined).

    Handles BOTH actors' confirmed output shapes:
    - brilliant_gum: "images" (array of URLs), "video" (object with
      stream_url), "reshared" / "reshared_from"
    - data-slayer:   "attachments" (array of typed objects),
      "is_repost" / "shared_post"
    """
    # ── HarvestAPI shape (CONFIRMED from official docs:
    #    postImages / postVideo / article / newsletterUrl / repost) ──
    #    For pure reposts with no media of their own, classify by the
    #    reposted content's media instead.
    src = post
    repost_obj = post.get("repost")
    if isinstance(repost_obj, dict) and repost_obj:
        has_own_media = (post.get("postImages") or post.get("postVideo") or post.get("article"))
        if not has_own_media:
            src = repost_obj

    post_video = src.get("postVideo")
    if isinstance(post_video, dict) and (post_video.get("videoUrl") or post_video.get("thumbnailUrl")):
        return "videos"
    post_images = src.get("postImages")
    if isinstance(post_images, list) and len(post_images) > 0:
        return "multi_image" if len(post_images) >= 2 else "single_image"
    document_obj = src.get("document")
    if isinstance(document_obj, dict) and document_obj:
        return "documents"
    if src.get("newsletterUrl") or (isinstance(src.get("article"), dict) and src.get("article")):
        return "other"          # Articles & Newsletters

    # ── brilliant_gum shape (images array / video object) ───────────
    video_obj = post.get("video")
    if isinstance(video_obj, dict) and video_obj.get("stream_url"):
        return "videos"
    images = post.get("images")
    if isinstance(images, list):
        if len(images) >= 2:
            return "multi_image"
        if len(images) == 1:
            return "single_image"

    # ── data-slayer shape (attachments array) ───────────────────────
    attachments = post.get("attachments") or []
    if not attachments and post.get("is_repost") and isinstance(post.get("shared_post"), dict):
        attachments = post["shared_post"].get("attachments") or []

    types = [str(a.get("type","")).lower() for a in attachments if isinstance(a, dict)]

    if any("video" in t for t in types):
        return "videos"
    if any("document" in t or "pdf" in t for t in types):
        return "documents"
    image_count = sum(1 for t in types if "image" in t or "photo" in t)
    if image_count >= 2:
        return "multi_image"
    if image_count == 1:
        return "single_image"
    if post.get("poll") or _get_first(post, ["poll_options","pollOptions"]):
        return "polls"

    text = str(post.get("text","") or post.get("commentary","") or post.get("content","")).lower()
    if not attachments and not text:
        return "text_only"
    if "linkedin.com/pulse" in text or _get_first(post, ["article_url","articleUrl"]):
        return "other"          # Article
    if any(k in text for k in ["newsletter", "subscribe to our"]):
        return "other"          # Newsletter
    if not attachments:
        return "text_only"
    return "other"               # Celebrations / anything unrecognized


def fetch_linkedin(company_handle, brand_name=None):
    print(f"💼 Fetching LinkedIn: {company_handle}")
    company_name = company_handle.replace("-"," ").replace("_"," ")
    data = {
        "company":   company_handle,
        "followers": "N/A",
        "employees": "N/A",
        "snippet":   "",
        "url":       f"https://www.linkedin.com/company/{company_handle}",
    }

    def extract_followers(text):
        for pattern in [
            r'([\d,]+\.?\d*[KMB]?\+?)\s*followers',
            r'followers[:\s]+([\d,]+\.?\d*[KMB]?\+?)',
            r'([\d,]+)\s*follow',
        ]:
            m = re.search(pattern, text, re.I)
            if m:
                return m.group(1).strip()
        return None

    def extract_employees(text):
        for pattern in [
            r'([\d,\-]+\.?\d*[KMB]?\+?)\s*employees',
            r'employees[:\s]+([\d,\-]+)',
            r'(\d+[-–]\d+)\s*employees',
        ]:
            m = re.search(pattern, text, re.I)
            if m:
                return m.group(1).strip()
        return None

    queries = [f'{company_name} linkedin followers']
    if brand_name and brand_name.lower() != company_name.lower():
        queries.append(f'{brand_name} linkedin followers')
        queries.append(f'"{brand_name}" linkedin company followers')
    queries.append(f'site:linkedin.com/company/{company_handle}')
    queries.append(f'{company_name} site:linkedin.com followers')

    for q in queries:
        try:
            print(f"   🔍 LinkedIn query: {q}")
            r = requests.get(
                "https://www.searchapi.io/api/v1/search",
                params={"engine":"google","q":q,
                        "api_key":SEARCHAPI_KEY,"num":5}
            )
            if r.status_code != 200:
                continue

            result_json = r.json()

            for key in ["knowledge_graph","answer_box"]:
                section = result_json.get(key, {})
                if section:
                    s = json.dumps(section)
                    f = extract_followers(s)
                    e = extract_employees(s)
                    if f and data["followers"] == "N/A":
                        data["followers"] = f
                        print(f"   ✅ {key} followers: {f}")
                    if e and data["employees"] == "N/A":
                        data["employees"] = e

            for res in result_json.get("organic_results", []):
                link     = res.get("link","")
                snippet  = res.get("snippet","")
                title    = res.get("title","")
                all_text = f"{title} {snippet}"
                sitelinks = res.get("sitelinks", [])
                if isinstance(sitelinks, list):
                    for sl in sitelinks:
                        all_text += f" {sl.get('title','')} {sl.get('snippet','')}"
                elif isinstance(sitelinks, dict):
                    all_text += json.dumps(sitelinks)

                f = extract_followers(all_text)
                e = extract_employees(all_text)
                if f and data["followers"] == "N/A":
                    data["followers"] = f
                    print(f"   ✅ Followers: {f}")
                if e and data["employees"] == "N/A":
                    data["employees"] = e
                if "linkedin.com/company" in link:
                    if not data.get("snippet"):
                        data["snippet"] = snippet
                    if title:
                        data["company"] = title.replace("| LinkedIn","").replace("LinkedIn","").strip()

            if data["followers"] != "N/A":
                print(f"   ✅ LinkedIn final: {data['followers']} followers")
                break  # stop trying more search queries, but still fall through
                       # to the post-level fetch below — a plain 'return' here
                       # would have skipped it entirely

        except Exception as e:
            print(f"   ⚠️ LinkedIn query failed: {e}")

    # ── Posts via Apify — reuses the exact same metric machinery as
    #    Facebook (day distribution, hour insights, consistency,
    #    momentum, per-category group metrics). ─────────────────────
    try:
        followers_n = parse_num(data.get("followers", 0))
        posts_raw = fetch_linkedin_posts_via_apify(company_handle)
        print(f"   🔍 LinkedIn posts fetched: {len(posts_raw)}")

        parsed = []
        for post in posts_raw:
            # Covers all three actors seen so far:
            # - harvestapi:    possibly nested "engagement" {likes, comments,
            #                  shares} and "postedAt" {timestamp/date}; text
            #                  in "content"
            # - brilliant_gum: num_likes / num_comments / num_reposts / posted
            # - data-slayer:   likes / comments / shares / created_at
            engagement = post.get("engagement") if isinstance(post.get("engagement"), dict) else {}
            likes    = _get_first(engagement, ["likes", "reactions"], None)
            if likes is None:
                likes = _get_first(post, ["num_likes", "likes", "reactions", "likesCount", "reactionsCount"], 0) or 0
            comments = _get_first(engagement, ["comments"], None)
            if comments is None:
                comments = _get_first(post, ["num_comments", "comments", "commentsCount"], 0) or 0
            shares   = _get_first(engagement, ["shares", "reposts"], None)
            if shares is None:
                shares = _get_first(post, ["num_reposts", "shares", "reposts", "sharesCount", "repostsCount"])

            url = _get_first(post, ["linkedinUrl", "url", "postUrl", "link"]) or "N/A"

            # Timestamp: try each candidate until one actually PARSES —
            # a plain "first non-empty" pick would grab relative strings
            # like "2 days ago" and silently lose the date. Nested
            # postedAt objects (harvestapi) are checked first.
            ts = None
            posted_at = post.get("postedAt")
            if isinstance(posted_at, dict):
                for k in ["timestamp", "date", "postedDate"]:
                    candidate = posted_at.get(k)
                    if candidate not in (None, ""):
                        ts = _to_timestamp(candidate)
                        if ts is not None:
                            break
            if ts is None:
                for ts_field in ["posted", "postedAt", "created_at", "timestamp", "date", "publishedAt", "postedDate", "time"]:
                    candidate = post.get(ts_field)
                    if isinstance(candidate, dict):
                        continue
                    if candidate not in (None, ""):
                        ts = _to_timestamp(candidate)
                        if ts is not None:
                            break

            ptype    = _classify_linkedin_post(post)
            try: likes = int(likes)
            except: likes = 0
            try: comments = int(comments)
            except: comments = 0
            parsed.append({
                "type": ptype, "likes": likes, "comments": comments,
                "shares": (int(shares) if shares is not None else None),
                "reposts": None, "views": None, "url": url, "ts": ts,
            })

        posts_with_dates    = [pp for pp in parsed if pp["ts"] is not None]
        posts_without_dates = [pp for pp in parsed if pp["ts"] is None]
        posts_with_dates.sort(key=lambda pp: pp["ts"], reverse=True)
        parsed = (posts_with_dates + posts_without_dates)[:30]

        post_count = len(parsed)
        li_categories = {}
        for cat in ["videos", "multi_image", "single_image", "documents", "polls", "text_only", "other"]:
            group = [pp for pp in parsed if pp["type"] == cat]
            li_categories[cat] = _compute_group_metrics(group, followers_n)
            li_categories[cat]["n"] = len(group)

        total_l = sum(pp["likes"] for pp in parsed)
        total_c = sum(pp["comments"] for pp in parsed)
        engagement_total = total_l + total_c
        eng_rate = f"{(engagement_total/post_count/followers_n*100):.2f}%" if (post_count and followers_n) else "N/A"

        timestamps = [pp["ts"] for pp in parsed if pp["ts"]]
        if len(timestamps) >= 2:
            span_days = max((max(timestamps) - min(timestamps)) / 86400, 1)
            posting_frequency = f"{(post_count / span_days * 7):.1f} / week"
        else:
            posting_frequency = "N/A"

        day_dist      = compute_day_distribution(parsed)
        hour_insights = compute_hour_insights(parsed)
        consistency   = compute_posting_consistency(parsed)
        momentum      = compute_momentum(parsed)

        data.update({
            "sample_size":       post_count,
            "engagement_rate":   eng_rate,
            "engagement_total":  engagement_total,
            "posting_frequency": posting_frequency,
            "most_active_day":  day_dist["most_active_day"],
            "least_active_day": day_dist["least_active_day"],
            "best_performing_day": day_dist["best_performing_day"],
            "avg_engagement": round(engagement_total / post_count, 1) if post_count else "N/A",
            "most_frequent_hour":   hour_insights["most_frequent_hour"],
            "best_performing_hour": hour_insights["best_performing_hour"],
            "outlier_note":         hour_insights["outlier_note"],
            "posting_consistency":  consistency,
            "momentum_pct":         momentum["momentum_pct"],
            "momentum_direction":   momentum["momentum_direction"],
            "categories": li_categories,
        })
        print(f"   ✅ {post_count} posts across 7 categories | ER: {eng_rate}")
    except Exception as e:
        import traceback
        print(f"   ⚠️ LinkedIn post-level fetch failed: {e}")
        print(traceback.format_exc())

    return data


# ═══════════════════════════════════════════════════════════════
# STEP 3 — CLAUDE ANALYSIS
# ═══════════════════════════════════════════════════════════════

def analyse_with_claude(website_url, handles, ig, fb, yt, li):
    print("\n🤖 Sending all data to Claude...")
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    prompt = f"""
You are a social media analyst writing slide copy for a branded report. Analyse this
brand's social media presence.

Website: {website_url}

INSTAGRAM DATA (overall stats plus per-format breakdown for Images, Carousels, Reels):
{json.dumps(ig, indent=2)}

FACEBOOK DATA:
{json.dumps(fb, indent=2)}

YOUTUBE DATA:
{json.dumps(yt, indent=2)}

LINKEDIN DATA:
{json.dumps(li, indent=2)}

WRITING RULES (apply to EVERY *_analysis / *_summary field):
- 2-3 sentences, strict maximum ~420 characters (fixed-size slide cards).
- EXCEPTIONS with SMALLER boxes — hard character limits:
  * youtube_content_analysis, youtube_comparison_analysis: max 2 sentences, ~250 characters
  * linkedin_content_analysis, linkedin_comparison_analysis: max 2 sentences, ~200 characters
  * every value inside cross_platform (strongest_platform, weakest_platform,
    content_consistency, growth_opportunity): 1 short sentence, ~120 characters
  * cross_platform.overall_recommendation: max ~300 characters
  * every "strength" and "recommendation" field: 1 short sentence, ~120 characters
- Use the REAL numbers from the data. Never invent numbers. Never write N/A if data exists.
- If a platform's sample_size is 0 or missing, say so plainly instead of inventing analysis.
- Benchmark bands for Engagement Rate per follower: under 1% = Low; 1-3.5% = Average/healthy;
  3.5-6% = Good; 6%+ = Excellent. Compare rates against the brand's industry/category norms
  where you know them. YouTube uses engagement-by-VIEW: ~2% average; 5%+ typical for Shorts;
  2-5% typical for long-form.
- "er_estimated" fields already include a 10% uplift for shares/saves — use the estimated
  rate when analysing Images/Carousels/Photos, but do NOT explain how it is calculated.
- Overview analyses must cover: posting_frequency vs best practice (IG/FB 3-5/week,
  LinkedIn 2-3/week), posting_consistency, a suggestion to post at best_performing_hour,
  and the video-format engagement vs the category's market standard.
- Content-strategy analyses must say what the brand posts MOST and whether that matches
  where the engagement is. For LinkedIn, ignore the "other" category in this judgement.
- Comparison analyses must name which format works best (highest er_per_follower) and call
  out anything unusual vs standard behaviour (e.g. unusually high shares).
- Deep-dive analyses must compare engagement rate BY VIEW to the industry standard for this
  brand's category and give one concrete fix.

Return ONLY a JSON object with EXACTLY these keys:
{{
  "brand_name": "brand name",
  "niche": "3-5 words describing the brand",
  "quick_overview": "3-4 sentences combining all platforms: which platform is primary and why, one line on each other platform",
  "instagram": {{"strength": "biggest strength", "recommendation": "one actionable tip"}},
  "instagram_analysis": "overview analysis per rules",
  "instagram_content_analysis": "what they post most on Instagram",
  "instagram_images_analysis": "images performance, use er_estimated",
  "instagram_carousels_analysis": "carousels performance, use er_estimated",
  "instagram_reels_analysis": "deep-dive: er_by_view vs category standard + one fix",
  "instagram_comparison_analysis": "which of images/carousels/reels works best + unusual behaviour",
  "instagram_summary": "overall Instagram verdict incl. momentum",
  "facebook": {{"strength": "biggest strength", "recommendation": "one actionable tip"}},
  "facebook_analysis": "overview analysis per rules",
  "facebook_content_analysis": "what they post most on Facebook",
  "facebook_photos_analysis": "photos performance, use er_estimated",
  "facebook_videos_analysis": "deep-dive: er_by_view vs standard + one fix",
  "facebook_comparison_analysis": "which of photos/videos/links works best",
  "facebook_summary": "overall Facebook verdict incl. momentum",
  "youtube": {{"strength": "biggest strength", "recommendation": "one actionable tip"}},
  "youtube_analysis": "overview analysis per rules",
  "youtube_content_analysis": "what they post most: shorts vs long-form",
  "youtube_shorts_analysis": "deep-dive: shorts er_by_view vs standard + one fix",
  "youtube_videos_analysis": "deep-dive: long-form er_by_view vs standard + one fix",
  "youtube_comparison_analysis": "shorts vs videos: which wins on engagement by view vs reach, vs industry standard",
  "linkedin": {{"strength": "biggest strength", "recommendation": "one actionable tip"}},
  "linkedin_analysis": "overview analysis per rules",
  "linkedin_content_analysis": "what they post most (ignore 'other')",
  "linkedin_category_analyses": {{
    "videos": "1-2 sentences", "multi_image": "1-2 sentences", "single_image": "1-2 sentences",
    "documents": "1-2 sentences", "polls": "1-2 sentences", "text_only": "1-2 sentences",
    "other": "1-2 sentences"
  }},
  "linkedin_comparison_analysis": "which format has max er_per_follower — conclude it works best",
  "linkedin_summary": "overall LinkedIn verdict incl. momentum",
  "cross_platform": {{
    "strongest_platform": "which and why",
    "weakest_platform": "which and why",
    "content_consistency": "consistent across platforms?",
    "growth_opportunity": "biggest growth opportunity",
    "overall_recommendation": "top priority action"
  }},
  "cross_platform_summary": "3-4 sentences: all four platforms side by side, red flags included",
  "overall_analysis": "3-4 sentences framed as: overall it's good right now, and here is what can be improved"
}}
Return ONLY JSON. No markdown. No explanation.
"""
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        messages=[{"role":"user","content":prompt}]
    )
    raw = msg.content[0].text.strip()
    try:
        return json.loads(raw)
    except:
        return json.loads(raw.replace("```json","").replace("```","").strip())


# ═══════════════════════════════════════════════════════════════
# STEP 4 — MIDNIGHT PANDA BRANDED PPT HELPERS (REDESIGN)
# ═══════════════════════════════════════════════════════════════

# Logo files — put logo_white.png (for dark slides) and logo_black.png
# (for light slides) inside a "static" folder next to app.py. If the
# files are missing the brand mark falls back to the old dot design so
# report generation never fails.
LOGO_WHITE = os.path.join("static", "logo_white.png")
LOGO_BLACK = os.path.join("static", "logo_black.png")

def usfmt(n):
    """US comma format: 87981 -> '87,981'. Leaves non-numeric values alone."""
    try:
        return f"{int(str(n).replace(',','')):,}"
    except Exception:
        return str(n)

def bg(slide, prs, color):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def tb(slide, text, x, y, w, h, size=14, bold=False, color=TEXT_DARK,
       align=PP_ALIGN.LEFT, font=FONT_MONO, italic=False, line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.text = str(text); p.alignment = align
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.italic = italic
    p.font.name = font
    if line_spacing:
        p.line_spacing = line_spacing
    return box

def brand_mark(slide, dark_bg=False, top_right=True):
    """Panda logo + MIDNIGHT PANDA name, identical location on every slide.
    White logo on dark slides, black logo on light slides."""
    color = TEXT_LIGHT if dark_bg else TEXT_DARK
    logo  = LOGO_WHITE if dark_bg else LOGO_BLACK
    lx = 10.10 if top_right else 0.55
    if os.path.exists(logo):
        slide.shapes.add_picture(logo, Inches(lx), Inches(0.40), height=Inches(0.36))
        tb(slide, "MIDNIGHT PANDA", lx + 0.46, 0.44, 2.30, 0.32, size=9.5, color=color, font=FONT_MONO)
    else:
        # Fallback: old dot mark, so a missing logo file never breaks a report
        inner = TEXT_DARK if dark_bg else TEXT_LIGHT
        for ox in [lx, lx + 0.24]:
            c1 = slide.shapes.add_shape(9, Inches(ox), Inches(0.45), Inches(0.16), Inches(0.16))
            c1.fill.solid(); c1.fill.fore_color.rgb = color; c1.line.fill.background()
            c2 = slide.shapes.add_shape(9, Inches(ox+0.05), Inches(0.50), Inches(0.05), Inches(0.05))
            c2.fill.solid(); c2.fill.fore_color.rgb = inner; c2.line.fill.background()
        tb(slide, "MIDNIGHT PANDA", lx + 0.53, 0.42, 2.20, 0.32, size=9.5, color=color, font=FONT_MONO)

def footer_bar(slide, page_num, dark_bg=False):
    line_color = RGBColor(0x2A,0x2A,0x2A) if dark_bg else RGBColor(0xD8,0xD5,0xCC)
    ln = slide.shapes.add_shape(1, Inches(0), Inches(7.46), Inches(13.33), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = line_color; ln.line.fill.background()
    tb(slide, "MIDNIGHTPANDA.AI · SOCIAL INTELLIGENCE", 0.55, 7.08, 5.00, 0.30,
       size=8.5, color=TEXT_FOOTER, font=FONT_MONO)
    tb(slide, f"{page_num:02d} / {TOTAL_SLIDES}", 11.93, 7.08, 0.85, 0.30,
       size=8.5, color=TEXT_FOOTER, font=FONT_MONO, align=PP_ALIGN.RIGHT)

def kicker_header(slide, kicker, title, subtitle, dark_bg=False):
    title_color = TEXT_LIGHT if dark_bg else TEXT_DARK
    sub_color   = TEXT_GRAY_LT if dark_bg else TEXT_GRAY
    tb(slide, kicker.upper(), 0.55, 0.42, 6.00, 0.30, size=11, color=GOLD, font=FONT_MONO_MED)
    tb(slide, title, 0.55, 0.72, 11.50, 0.65, size=26, color=title_color, font=FONT_SERIF)
    if subtitle:
        tb(slide, subtitle, 0.55, 1.34, 11.50, 0.35, size=13, color=sub_color, font=FONT_MONO_LT)

def stat_block(slide, x, y, w, value, label, size=54, dark_bg=False):
    label_color = TEXT_GRAY_LT if dark_bg else TEXT_GRAY
    vh = 0.95 if size >= 50 else (0.75 if size >= 40 else 0.55)
    tb(slide, value, x, y, w, vh, size=size, color=GOLD, font=FONT_STAT)
    tb(slide, label.upper(), x, y+vh+0.05, w, 0.30, size=9.5, color=label_color, font=FONT_MONO)

def card(slide, x, y, w, h, label, body, dark_bg=False, body_size=12.5):
    """Card with AUTO-FIT body text: if the text is too long for the box,
    the font steps down (to a 9.5pt floor) and, if still too long, the
    text is cut at the last full sentence that fits. This guarantees text
    can never overflow a card, no matter how long Claude's analysis runs."""
    fill_color = CARD_DARK if dark_bg else CARD_LIGHT
    body_color = TEXT_LIGHT if dark_bg else TEXT_DARK
    r = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill_color; r.line.fill.background()
    try: r.adjustments[0] = 0.06
    except: pass
    tb(slide, label.upper(), x+0.28, y+0.18, w-0.56, 0.30, size=10.5, color=GOLD, font=FONT_MONO_MED)

    body = str(body or "")
    body_w = w - 0.56          # usable width in inches
    body_h = h - 0.62          # usable height (label + padding removed)

    def fits(text, size):
        cpl = max(int(body_w * 72 / (size * 0.60)), 1)     # chars per line
        import math as _m
        lines = sum(_m.ceil(max(len(ln), 1) / cpl) for ln in text.split("\n"))
        return lines * size * 1.30 / 72 <= body_h

    size = body_size
    while size > 9.5 and not fits(body, size):
        size -= 0.5
    if not fits(body, size):
        cpl = max(int(body_w * 72 / (size * 0.60)), 1)
        max_lines = max(int(body_h * 72 / (size * 1.30)), 1)
        body = truncate_to_sentence(body, max_chars=cpl * max_lines)

    tb(slide, body, x+0.28, y+0.50, w-0.56, h-0.6, size=size, color=body_color,
       font=FONT_MONO_LT, line_spacing=1.15)

def shorten_url(url, maxlen=46):
    if not url or url == "N/A":
        return "N/A"
    return url if len(url) <= maxlen else url[:maxlen-1] + "…"

def link_list_card(slide, x, y, w, h, label, items, dark_bg=False):
    """Stacked "label: link" lines. Links are gold (never blue — the file's
    theme hyperlink color is also forced to gold after save)."""
    fill_color = CARD_DARK if dark_bg else CARD_LIGHT
    body_color = TEXT_LIGHT if dark_bg else TEXT_DARK
    r = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill_color; r.line.fill.background()
    try: r.adjustments[0] = 0.06
    except: pass
    tb(slide, label.upper(), x+0.28, y+0.14, w-0.56, 0.22, size=10, color=GOLD, font=FONT_MONO_MED)
    n = max(len(items), 1)
    line_h = (h - 0.40) / n
    for i, (prefix, url) in enumerate(items):
        box = slide.shapes.add_textbox(Inches(x+0.28), Inches(y+0.38+i*line_h), Inches(w-0.56), Inches(line_h))
        tf = box.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = prefix
        r1.font.size = Pt(11.5); r1.font.name = FONT_MONO_LT; r1.font.color.rgb = body_color
        r2 = p.add_run()
        if url and url != "N/A":
            r2.text = shorten_url(url)
            r2.hyperlink.address = url
        else:
            r2.text = "N/A"
        r2.font.color.rgb = GOLD
        r2.font.size = Pt(11.5); r2.font.name = FONT_MONO

def rounded_table(slide, headers, rows, x, y, w, dark_bg=False, row_h=0.52):
    """Curved-edge table built from rounded-rectangle rows."""
    cols = len(headers)
    cw = w / cols
    hdr_fill  = GOLD if dark_bg else BG_DARK
    hdr_text  = BG_DARK if dark_bg else TEXT_LIGHT
    body_fill = CARD_DARK if dark_bg else CARD_LIGHT
    body_text = TEXT_LIGHT if dark_bg else TEXT_DARK
    r = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(row_h))
    r.fill.solid(); r.fill.fore_color.rgb = hdr_fill; r.line.fill.background()
    try: r.adjustments[0] = 0.28
    except: pass
    for ci, hdr in enumerate(headers):
        tb(slide, hdr, x+ci*cw+0.18, y+0.10, cw-0.30, 0.32, size=12.5, bold=True,
           color=hdr_text, font=FONT_MONO_MED,
           align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    yy = y + row_h + 0.10
    for row in rows:
        rr = slide.shapes.add_shape(5, Inches(x), Inches(yy), Inches(w), Inches(row_h))
        rr.fill.solid(); rr.fill.fore_color.rgb = body_fill; rr.line.fill.background()
        try: rr.adjustments[0] = 0.28
        except: pass
        for ci, val in enumerate(row):
            tb(slide, str(val), x+ci*cw+0.18, yy+0.10, cw-0.30, 0.32, size=12,
               color=body_text, font=FONT_MONO,
               align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
        yy += row_h + 0.10
    return yy

def start_slide(prs, blank, slide_num):
    is_dark = (slide_num % 2 == 1)
    s = prs.slides.add_slide(blank)
    bg(s, prs, BG_DARK if is_dark else BG_LIGHT)
    brand_mark(s, dark_bg=is_dark, top_right=(slide_num != 1))
    return s, is_dark

def truncate_to_sentence(text, max_chars=480):
    """Hard cap for card text — cuts at the last full sentence that fits."""
    if not text or len(text) <= max_chars:
        return text
    window = text[:max_chars]
    best_cut = -1
    for sep in [". ", "! ", "? "]:
        cut = window.rfind(sep)
        if cut > best_cut:
            best_cut = cut
    if best_cut > 0:
        return window[:best_cut+1]
    cut = window.rfind(" ")
    return (window[:cut] if cut > 0 else window) + "…"

def _force_gold_hyperlinks(pptx_path):
    """Post-save: rewrite the theme's hyperlink colors to brand gold so
    links never render blue in PowerPoint / Keynote / Google Slides."""
    import zipfile, shutil, tempfile
    tmp = tempfile.mkdtemp()
    try:
        with zipfile.ZipFile(pptx_path) as z:
            z.extractall(tmp)
        theme_path = os.path.join(tmp, "ppt", "theme", "theme1.xml")
        if os.path.exists(theme_path):
            xml = open(theme_path, encoding="utf-8").read()
            xml = re.sub(r'(<a:hlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}(")', r'\g<1>B8945A\g<2>', xml)
            xml = re.sub(r'(<a:folHlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}(")', r'\g<1>B8945A\g<2>', xml)
            open(theme_path, "w", encoding="utf-8").write(xml)
            os.remove(pptx_path)
            with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as z:
                for root, dirs, files in os.walk(tmp):
                    for f in files:
                        fp = os.path.join(root, f)
                        z.write(fp, os.path.relpath(fp, tmp))
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ═══════════════════════════════════════════════════════════════
# STEP 5 — BUILD 40-SLIDE MIDNIGHT PANDA BRANDED PPT (REDESIGN)
# ═══════════════════════════════════════════════════════════════

def create_ppt(analysis, handles, ig_raw, fb_raw, yt_raw, li_raw, website_url):
    prs   = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    A = analysis  # shorthand
    ig = A.get("instagram", {})
    fb = A.get("facebook", {})
    yt = A.get("youtube", {})
    li = A.get("linkedin", {})
    cp = A.get("cross_platform", {})
    brand = A.get("brand_name", website_url)
    niche = A.get("niche", "")
    site_display = website_url.replace("https://","").replace("http://","")

    def T(key, max_chars=480):
        """Analysis text with hard overflow cap."""
        return truncate_to_sentence(str(A.get(key, "") or ""), max_chars=max_chars)

    ig_images    = ig_raw.get("images", {})
    ig_carousels = ig_raw.get("carousels", {})
    ig_reels     = ig_raw.get("reels", {})
    yt_shorts    = yt_raw.get("shorts", {})
    yt_videos    = yt_raw.get("videos_performance", {})
    fb_photos    = fb_raw.get("photos", {})
    fb_videos    = fb_raw.get("videos", {})
    fb_links     = fb_raw.get("links", {})
    li_categories = li_raw.get("categories", {})

    # ── SLIDE 1: COVER (dark) ─────────────────────────────────
    s, dark = start_slide(prs, blank, 1)
    tb(s, "SOCIAL INTELLIGENCE REPORT", 0.55, 1.30, 6.00, 0.30, size=11, color=GOLD, font=FONT_MONO_MED)
    tb(s, brand, 0.55, 1.70, 11.50, 1.50, size=58, color=TEXT_LIGHT, font=FONT_SERIF)
    tb(s, niche, 0.55, 2.70, 11.50, 0.45, size=15, color=TEXT_GRAY_LT, font=FONT_MONO_LT)
    tb(s, site_display, 0.55, 3.17, 11.50, 0.35, size=11.5, color=GOLD, font=FONT_MONO)
    box = s.shapes.add_textbox(Inches(0.55), Inches(3.70), Inches(8.0), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = "Generated by "
    r1.font.size = Pt(12); r1.font.name = FONT_MONO_LT; r1.font.color.rgb = TEXT_GRAY_LT
    r2 = p.add_run(); r2.text = "midnightpanda.ai"
    r2.hyperlink.address = "https://midnightpanda.ai"
    r2.font.size = Pt(12); r2.font.name = FONT_MONO_MED; r2.font.color.rgb = GOLD
    tb(s, "OUR OTHER TOOLS", 0.55, 4.55, 6.00, 0.30, size=11, color=GOLD, font=FONT_MONO_MED)
    for i, tool in enumerate(["AI Brand Intelligence", "Competitor Ad Analysis",
                              "Social Listening", "SEO & Traffic Intelligence"]):
        x = 0.55 + (i % 2) * 4.55
        y = 4.95 + (i // 2) * 0.78
        r = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(4.25), Inches(0.60))
        r.fill.solid(); r.fill.fore_color.rgb = CARD_DARK; r.line.fill.background()
        try: r.adjustments[0] = 0.30
        except: pass
        tb(s, tool, x+0.30, y+0.14, 3.85, 0.34, size=12.5, color=TEXT_LIGHT, font=FONT_MONO)
    ln = s.shapes.add_shape(1, Inches(0), Inches(7.46), Inches(13.33), Pt(0.75))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x2A,0x2A,0x2A); ln.line.fill.background()
    tb(s, f"{TOTAL_SLIDES}-slide report  ·  Instagram · Facebook · YouTube · LinkedIn",
       0.55, 7.08, 8.00, 0.30, size=8.5, color=TEXT_FOOTER, font=FONT_MONO)

    # ── SLIDE 2: Report Overview (light) ──────────────────────
    s, dark = start_slide(prs, blank, 2)
    kicker_header(s, "Report Overview", "Social Media Overview", f"{brand} — all platforms at a glance", dark_bg=dark)
    rounded_table(s, ["Platform", "Handle", "Followers", "Status"], [
        ["Instagram", f"@{handles.get('instagram','N/A')}", usfmt(ig_raw.get("followers","N/A")),
         "Active" if handles.get("instagram") else "Not Found"],
        ["Facebook",  handles.get("facebook","N/A"),  usfmt(fb_raw.get("followers","N/A")),
         "Active" if handles.get("facebook") else "Not Found"],
        ["YouTube",   yt_raw.get("handle", handles.get("youtube","N/A")), usfmt(yt_raw.get("subscribers","N/A")),
         "Active" if handles.get("youtube") else "Not Found"],
        ["LinkedIn",  handles.get("linkedin","N/A"),  usfmt(li_raw.get("followers","N/A")),
         "Active" if handles.get("linkedin") else "Not Found"],
    ], 0.55, 1.70, 12.23, dark_bg=dark)
    card(s, 0.55, 4.95, 12.23, 1.95, "Quick Overview of Report", T("quick_overview"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 2, dark_bg=dark)

    # ── Shared: hours/days/consistency/momentum rows ───────────
    def hours_days_rows(slide, d, y, dark_bg):
        sub = TEXT_GRAY_LT if dark_bg else TEXT_GRAY
        tb(slide, f"Most Frequent Hour: {d.get('most_frequent_hour','N/A')}   ·   Best Performing Hour: {d.get('best_performing_hour','N/A')}",
           0.55, y, 12.23, 0.24, size=12, color=sub, font=FONT_MONO_LT)
        tb(slide, f"Most Active Day: {d.get('most_active_day','N/A')}   ·   Least Active Day: {d.get('least_active_day','N/A')}   ·   Best Performing Day: {d.get('best_performing_day','N/A')}",
           0.55, y+0.27, 12.23, 0.24, size=12, color=sub, font=FONT_MONO_LT)
        tb(slide, f"Posting Consistency: {d.get('posting_consistency','N/A')}",
           0.55, y+0.54, 12.23, 0.24, size=12, color=sub, font=FONT_MONO_LT)
        # Momentum on its own line BELOW consistency — always brand gold
        md = d.get("momentum_direction")
        arrow = "\u25B2 " if md == "up" else ("\u25BC " if md == "down" else "")
        tb(slide, f"Momentum: {arrow}{d.get('momentum_pct','N/A')}", 0.55, y+0.81, 12.23, 0.24,
           size=12, bold=True, color=GOLD, font=FONT_MONO_MED)

    def overview_slide(n, kicker, title, sub, stats_row1, stats_row2, d, analysis_text):
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, kicker, title, sub, dark_bg=dark)
        for i, (val, lbl, sz) in enumerate(stats_row1):
            stat_block(s, 0.55 + i*4.18, 1.75, 3.86, val, lbl, size=sz, dark_bg=dark)
        for i, (val, lbl, sz) in enumerate(stats_row2):
            stat_block(s, 0.55 + i*4.18, 3.05, 3.86, val, lbl, size=sz, dark_bg=dark)
        div = RGBColor(0x2A,0x2A,0x2A) if dark else RGBColor(0xD8,0xD5,0xCC)
        ln = s.shapes.add_shape(1, Inches(0.55), Inches(4.02), Inches(12.23), Pt(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = div; ln.line.fill.background()
        hours_days_rows(s, d, 4.14, dark)
        card(s, 0.55, 5.35, 12.23, 1.60, f"{kicker} Analysis", analysis_text, dark_bg=dark, body_size=11)
        footer_bar(s, n, dark_bg=dark)

    # ── SLIDE 3: Instagram Overview (dark) ─────────────────────
    overview_slide(3, "Instagram", "Instagram Overview",
        f"@{handles.get('instagram','N/A')}  ·  Based on last {ig_raw.get('sample_size','N/A')} posts",
        [(usfmt(ig_raw.get("followers","N/A")), "Followers", 54),
         (usfmt(ig_raw.get("posts","N/A")), "Total Posts (all time)", 54),
         (ig_raw.get("posting_frequency","N/A"), "Posting Frequency", 44)],
        [(ig_raw.get("engagement_rate","N/A"), "Engagement Rate / Follower", 32),
         (ig_raw.get("engagement_rate_reels","N/A"), "Engagement Rate / Views (Reels Only)", 32),
         (ig_raw.get("avg_engagement","N/A"), "Avg Engagement (Total ÷ Posts)", 32)],
        ig_raw, T("instagram_analysis"))

    # ── SLIDE 4: IG Content Strategy (light) ───────────────────
    s, dark = start_slide(prs, blank, 4)
    kicker_header(s, "Instagram", "Content Strategy",
                  f"Content mix from the last {ig_raw.get('sample_size',30)} posts", dark_bg=dark)
    img_c = ig_raw.get("img_count",0); car_c = ig_raw.get("car_count",0); vid_c = ig_raw.get("vid_count",0)
    total = max(img_c + car_c + vid_c, 1)
    cd = ChartData(); cd.categories = ["Images","Carousels","Videos / Reels"]
    cd.add_series("Posts",(img_c,car_c,vid_c))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.85),Inches(1.80),Inches(4.90),Inches(3.60),cd).chart
    chart.has_legend = True
    for i, col in enumerate([PIE_TAN, CARD_DARK, GOLD]):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    for lbl, cnt, y in [("Videos / Reels",vid_c,1.95),("Carousels",car_c,3.05),("Images",img_c,4.15)]:
        stat_block(s, 6.55, y, 5.60, f"{int(round(cnt/total*100))}%", f"{lbl} — {cnt} posts", size=32, dark_bg=dark)
    card(s, 0.55, 5.60, 12.23, 1.30, "Content Strategy Analysis", T("instagram_content_analysis"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 4, dark_bg=dark)

    # ── Format performance slide (IG images/carousels, FB photos) ──
    def format_perf_slide(n, kicker, title, d, sample, unit, analysis_text):
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, kicker, title,
                      f"Based on {d.get('n',0)} {unit} from the last {sample}", dark_bg=dark)
        stat_block(s, 0.55, 1.75, 3.86, d.get("avg_engagement","N/A"), "Avg Engagement / Post", size=44, dark_bg=dark)
        stat_block(s, 4.73, 1.75, 3.86, d.get("er_per_follower","N/A"), "Engagement Rate / Followers", size=44, dark_bg=dark)
        stat_block(s, 8.92, 1.75, 3.86, d.get("er_estimated","N/A"), "Estimated Engagement Rate / Followers", size=44, dark_bg=dark)
        stat_block(s, 0.55, 3.10, 3.86, d.get("avg_likes","N/A"), "Avg Likes / Post", size=32, dark_bg=dark)
        stat_block(s, 4.73, 3.10, 3.86, d.get("avg_comments","N/A"), "Avg Comments / Post", size=32, dark_bg=dark)
        stat_block(s, 8.92, 3.10, 3.86, d.get("n",0), "Posts in Sample", size=32, dark_bg=dark)
        card(s, 0.55, 4.45, 12.23, 2.10, "Analysis", analysis_text, dark_bg=dark)
        footer_bar(s, n, dark_bg=dark)

    format_perf_slide(5, "Instagram — Images", "Images Performance", ig_images,
                      ig_raw.get("sample_size",30), "image posts", T("instagram_images_analysis"))
    format_perf_slide(6, "Instagram — Carousels", "Carousels Performance", ig_carousels,
                      ig_raw.get("sample_size",30), "carousel posts", T("instagram_carousels_analysis"))

    # ── SLIDE 7: IG Reels Performance (dark) — ER by view on top ──
    s, dark = start_slide(prs, blank, 7)
    kicker_header(s, "Instagram — Reels", "Reels Performance",
                  f"Based on {ig_reels.get('n',0)} reels from the last {ig_raw.get('sample_size',30)}", dark_bg=dark)
    stat_block(s, 0.55, 1.75, 3.86, ig_reels.get("er_by_view","N/A"), "Engagement Rate by View", size=44, dark_bg=dark)
    stat_block(s, 4.73, 1.75, 3.86, ig_reels.get("avg_engagement","N/A"), "Avg Engagement / Reel", size=44, dark_bg=dark)
    stat_block(s, 8.92, 1.75, 3.86, ig_reels.get("avg_views","N/A"), "Avg Views / Reel", size=44, dark_bg=dark)
    stat_block(s, 0.55, 3.10, 3.86, ig_reels.get("avg_likes","N/A"), "Avg Likes / Reel", size=28, dark_bg=dark)
    stat_block(s, 4.73, 3.10, 3.86, ig_reels.get("avg_comments","N/A"), "Avg Comments / Reel", size=28, dark_bg=dark)
    stat_block(s, 8.92, 3.10, 3.86, ig_reels.get("er_per_follower","N/A"), "Engagement Rate / Followers", size=28, dark_bg=dark)
    stat_block(s, 0.55, 4.20, 3.86, ig_reels.get("like_rate","N/A"), "Like Rate (per view)", size=28, dark_bg=dark)
    stat_block(s, 4.73, 4.20, 3.86, ig_reels.get("comment_rate","N/A"), "Comment Rate (per view)", size=28, dark_bg=dark)
    stat_block(s, 8.92, 4.20, 3.86, ig_reels.get("share_rate","N/A"), "Share Rate (per view)", size=28, dark_bg=dark)
    stat_block(s, 0.55, 5.30, 3.86, ig_reels.get("er_estimated","N/A"), "Estimated ER / Followers", size=28, dark_bg=dark)
    stat_block(s, 4.73, 5.30, 3.86, ig_reels.get("n",0), "Reels in Sample", size=28, dark_bg=dark)
    footer_bar(s, 7, dark_bg=dark)

    # ── Deep dive slide (links + analysis) ─────────────────────
    def deep_dive_slide(n, kicker, title, d, analysis_text, unit="post"):
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, kicker, title,
                      "Top and bottom performers, and how the rate compares to the market", dark_bg=dark)
        items = [
            (f"Top {unit} (score {d.get('top_score','N/A')}):  ", d.get("top_url","N/A")),
            (f"Worst {unit} (score {d.get('worst_score','N/A')}):  ", d.get("worst_url","N/A")),
        ]
        if d.get("max_views") is not None:
            items.append((f"Max views ({usfmt(d.get('max_views','N/A'))}):  ", d.get("max_views_url","N/A")))
        link_list_card(s, 0.55, 1.85, 12.23, 1.60, "Top / Worst / Max Views", items, dark_bg=dark)
        card(s, 0.55, 3.80, 12.23, 2.60, "Analysis", analysis_text, dark_bg=dark)
        footer_bar(s, n, dark_bg=dark)

    # ── SLIDE 8: IG Reels Deep Dive (light) ────────────────────
    deep_dive_slide(8, "Instagram — Reels", "Reels Deep Dive", ig_reels,
                    T("instagram_reels_analysis"), unit="reel")

    # ── SLIDE 9: IG Format Comparison (dark) ───────────────────
    s, dark = start_slide(prs, blank, 9)
    kicker_header(s, "Instagram", "Format Comparison",
                  "Engagement rate per follower across the three formats", dark_bg=dark)
    rounded_table(s, ["Format", "Posts", "ER / Followers", "Estimated ER"], [
        ["Reels", ig_reels.get("n",0), ig_reels.get("er_per_follower","N/A"), ig_reels.get("er_estimated","N/A")],
        ["Carousels", ig_carousels.get("n",0), ig_carousels.get("er_per_follower","N/A"), ig_carousels.get("er_estimated","N/A")],
        ["Images", ig_images.get("n",0), ig_images.get("er_per_follower","N/A"), ig_images.get("er_estimated","N/A")],
    ], 0.55, 1.75, 12.23, dark_bg=dark)
    card(s, 0.55, 4.55, 12.23, 2.30, "What Works Best", T("instagram_comparison_analysis"), dark_bg=dark)
    footer_bar(s, 9, dark_bg=dark)

    # ── Summary slide (4 cards + verdict) ──────────────────────
    def summary_slide(n, kicker, title, cards_list):
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, kicker, title,
                      "Everything on one slide — metrics, momentum, and the verdict", dark_bg=dark)
        positions = [(0.55, 1.75), (6.83, 1.75), (0.55, 3.45), (6.83, 3.45)]
        for (x, y), (lbl, body) in zip(positions, cards_list[:4]):
            card(s, x, y, 5.96, 1.45, lbl, body, dark_bg=dark, body_size=11.5)
        if len(cards_list) > 4:
            lbl, body = cards_list[4]
            card(s, 0.55, 5.15, 12.23, 1.75, lbl, body, dark_bg=dark, body_size=11.5)
        footer_bar(s, n, dark_bg=dark)

    # ── SLIDE 10: IG Summary (light) ───────────────────────────
    summary_slide(10, "Instagram", "Instagram Summary", [
        ("Audience", f"{usfmt(ig_raw.get('followers','N/A'))} followers · primary reach platform"),
        ("Performance", f"ER {ig_raw.get('engagement_rate','N/A')} per follower · Reels {ig_raw.get('engagement_rate_reels','N/A')} per view · momentum {ig_raw.get('momentum_pct','N/A')}"),
        ("Cadence", f"{ig_raw.get('posting_frequency','N/A')} · {ig_raw.get('posting_consistency','N/A')} · best hour {ig_raw.get('best_performing_hour','N/A')} · best day {ig_raw.get('best_performing_day','N/A')}"),
        ("Strongest Format", str(ig.get("strength","N/A"))),
        ("Overall Instagram Analysis", T("instagram_summary")),
    ])

    # ── SLIDE 11: Facebook Overview (dark) ─────────────────────
    overview_slide(11, "Facebook", "Facebook Overview",
        f"{handles.get('facebook','N/A')}  ·  Based on last {fb_raw.get('sample_size','N/A')} posts",
        [(usfmt(fb_raw.get("followers","N/A")), "Page Followers", 54),
         (str(fb_raw.get("sample_size","N/A")), "Posts Analyzed", 54),
         (fb_raw.get("posting_frequency","N/A"), "Posting Frequency", 44)],
        [(fb_raw.get("engagement_rate","N/A"), "Engagement Rate / Follower", 32),
         (fb_raw.get("avg_engagement","N/A"), "Avg Engagement (Total ÷ Posts)", 32),
         (fb_raw.get("engagement_rate_videos","N/A"), "Engagement Rate (Videos Only)", 32)],
        fb_raw, T("facebook_analysis"))

    # ── SLIDE 12: FB Content Strategy (light) ──────────────────
    s, dark = start_slide(prs, blank, 12)
    kicker_header(s, "Facebook", "Content Strategy",
                  f"Content mix from the last {fb_raw.get('sample_size',30)} posts", dark_bg=dark)
    fb_photo_c = fb_raw.get("photo_count",0); fb_video_c = fb_raw.get("video_count",0); fb_link_c = fb_raw.get("link_count",0)
    fb_total = max(fb_photo_c + fb_video_c + fb_link_c, 1)
    cd = ChartData(); cd.categories = ["Photos","Videos","Links"]
    cd.add_series("Posts",(fb_photo_c,fb_video_c,fb_link_c))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.85),Inches(1.80),Inches(4.90),Inches(3.60),cd).chart
    chart.has_legend = True
    for i, col in enumerate([GOLD, CARD_DARK, PIE_TAN]):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    for lbl, cnt, y in [("Photos",fb_photo_c,1.95),("Videos",fb_video_c,3.05),("Links",fb_link_c,4.15)]:
        stat_block(s, 6.55, y, 5.60, f"{int(round(cnt/fb_total*100))}%", f"{lbl} — {cnt} posts", size=32, dark_bg=dark)
    card(s, 0.55, 5.60, 12.23, 1.30, "Content Strategy Analysis", T("facebook_content_analysis"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 12, dark_bg=dark)

    # ── SLIDE 13: FB Photos (dark) ─────────────────────────────
    format_perf_slide(13, "Facebook — Photos", "Photos Performance", fb_photos,
                      fb_raw.get("sample_size",30), "photo posts", T("facebook_photos_analysis"))

    # ── SLIDE 14: FB Videos Performance (light) ────────────────
    s, dark = start_slide(prs, blank, 14)
    kicker_header(s, "Facebook — Videos", "Videos Performance",
                  f"Based on {fb_videos.get('n',0)} videos from the last {fb_raw.get('sample_size',30)}", dark_bg=dark)
    stat_block(s, 0.55, 1.85, 3.86, fb_videos.get("er_by_view","N/A"), "Engagement Rate by View", size=44, dark_bg=dark)
    stat_block(s, 4.73, 1.85, 3.86, fb_videos.get("avg_engagement","N/A"), "Avg Engagement / Video", size=44, dark_bg=dark)
    stat_block(s, 8.92, 1.85, 3.86, fb_videos.get("avg_views","N/A"), "Avg Views / Video", size=44, dark_bg=dark)
    stat_block(s, 0.55, 3.45, 3.86, fb_videos.get("avg_likes","N/A"), "Avg Likes / Video", size=32, dark_bg=dark)
    stat_block(s, 4.73, 3.45, 3.86, fb_videos.get("avg_comments","N/A"), "Avg Comments / Video", size=32, dark_bg=dark)
    stat_block(s, 8.92, 3.45, 3.86, fb_videos.get("er_per_follower","N/A"), "Engagement Rate / Followers", size=32, dark_bg=dark)
    stat_block(s, 0.55, 5.00, 3.86, fb_videos.get("like_rate","N/A"), "Like Rate (per view)", size=32, dark_bg=dark)
    stat_block(s, 4.73, 5.00, 3.86, fb_videos.get("comment_rate","N/A"), "Comment Rate (per view)", size=32, dark_bg=dark)
    stat_block(s, 8.92, 5.00, 3.86, fb_videos.get("er_estimated","N/A"), "Estimated ER / Followers", size=32, dark_bg=dark)
    footer_bar(s, 14, dark_bg=dark)

    # ── SLIDE 15: FB Videos Deep Dive (dark) ───────────────────
    deep_dive_slide(15, "Facebook — Videos", "Videos Deep Dive", fb_videos,
                    T("facebook_videos_analysis"), unit="video")

    # ── SLIDE 16: FB Format Comparison (light) ─────────────────
    s, dark = start_slide(prs, blank, 16)
    kicker_header(s, "Facebook", "Format Comparison",
                  "Engagement rate per follower across the three formats", dark_bg=dark)
    rounded_table(s, ["Format", "Posts", "ER / Followers", "Estimated ER"], [
        ["Videos", fb_videos.get("n",0), fb_videos.get("er_per_follower","N/A"), fb_videos.get("er_estimated","N/A")],
        ["Photos", fb_photos.get("n",0), fb_photos.get("er_per_follower","N/A"), fb_photos.get("er_estimated","N/A")],
        ["Links",  fb_links.get("n",0),  fb_links.get("er_per_follower","N/A"),  fb_links.get("er_estimated","N/A")],
    ], 0.55, 1.75, 12.23, dark_bg=dark)
    card(s, 0.55, 4.55, 12.23, 2.30, "What Works Best", T("facebook_comparison_analysis"), dark_bg=dark)
    footer_bar(s, 16, dark_bg=dark)

    # ── SLIDE 17: FB Summary (dark) ────────────────────────────
    summary_slide(17, "Facebook", "Facebook Summary", [
        ("Audience", f"{usfmt(fb_raw.get('followers','N/A'))} followers"),
        ("Performance", f"Avg engagement {fb_raw.get('avg_engagement','N/A')} / post · videos {fb_videos.get('er_by_view','N/A')} ER by view · momentum {fb_raw.get('momentum_pct','N/A')}"),
        ("Cadence", f"{fb_raw.get('posting_frequency','N/A')} · {fb_raw.get('posting_consistency','N/A')} · best hour {fb_raw.get('best_performing_hour','N/A')} · best day {fb_raw.get('best_performing_day','N/A')}"),
        ("Strongest Format", str(fb.get("strength","N/A"))),
        ("Overall Facebook Analysis", T("facebook_summary")),
    ])

    # ── SLIDE 18: YouTube Overview (light) ─────────────────────
    overview_slide(18, "YouTube", "YouTube Overview",
        f"{yt_raw.get('handle', handles.get('youtube','N/A'))}  ·  Based on last {yt_raw.get('sample_size','N/A')} videos",
        [(usfmt(yt_raw.get("subscribers","N/A")), "Subscribers", 54),
         (str(yt_raw.get("sample_size","N/A")), "Videos Analyzed", 54),
         (yt_raw.get("posting_frequency","N/A"), "Posting Frequency", 44)],
        [(yt_raw.get("engagement_rate","N/A"), "Engagement Rate / Subscriber", 32),
         (yt_raw.get("avg_engagement","N/A"), "Avg Engagement (Total ÷ Videos)", 32),
         (yt_raw.get("avg_views_sample","N/A"), "Avg Views / Video (Total ÷ Videos)", 32)],
        yt_raw, T("youtube_analysis"))

    # ── SLIDE 19: YT Content Strategy (dark) ───────────────────
    s, dark = start_slide(prs, blank, 19)
    kicker_header(s, "YouTube", "Content Strategy",
                  f"Content mix from the last {yt_raw.get('sample_size',30)} videos", dark_bg=dark)
    yt_shorts_c = yt_raw.get("shorts_count",0); yt_videos_c = yt_raw.get("videos_count",0)
    yt_total = max(yt_shorts_c + yt_videos_c, 1)
    cd = ChartData(); cd.categories = ["Shorts","Videos"]
    cd.add_series("Videos",(yt_shorts_c, yt_videos_c))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.85),Inches(1.80),Inches(4.90),Inches(3.60),cd).chart
    chart.has_legend = True
    for i, col in enumerate([GOLD, RGBColor(0x3A,0x3A,0x3A)]):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    for lbl, cnt, y in [("Shorts",yt_shorts_c,2.30),("Videos",yt_videos_c,3.70)]:
        stat_block(s, 6.55, y, 5.60, f"{int(round(cnt/yt_total*100))}%",
                   f"{lbl} — {cnt} video{'s' if cnt != 1 else ''}", size=40, dark_bg=dark)
    card(s, 0.55, 5.60, 12.23, 1.30, "Content Strategy Analysis", T("youtube_content_analysis"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 19, dark_bg=dark)

    # ── YT stats slide (shared by Shorts & Videos) ─────────────
    def yt_stats_slide(n, kicker, title, d, sample, unit):
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, kicker, title,
                      f"Based on {d.get('n',0)} {unit} from the last {sample}", dark_bg=dark)
        stat_block(s, 0.55, 1.85, 3.86, d.get("er_by_view","N/A"), "Engagement Rate by View", size=44, dark_bg=dark)
        stat_block(s, 4.73, 1.85, 3.86, d.get("avg_engagement","N/A"), "Avg Engagement", size=44, dark_bg=dark)
        stat_block(s, 8.92, 1.85, 3.86, d.get("avg_views","N/A"), "Avg Views", size=44, dark_bg=dark)
        stat_block(s, 0.55, 3.45, 3.86, d.get("avg_likes","N/A"), "Avg Likes", size=32, dark_bg=dark)
        stat_block(s, 4.73, 3.45, 3.86, d.get("avg_comments","N/A"), "Avg Comments", size=32, dark_bg=dark)
        stat_block(s, 8.92, 3.45, 3.86, d.get("er_per_follower","N/A"), "Engagement Rate / Subscriber", size=32, dark_bg=dark)
        stat_block(s, 0.55, 5.00, 3.86, d.get("like_rate","N/A"), "Like Rate (per view)", size=32, dark_bg=dark)
        stat_block(s, 4.73, 5.00, 3.86, d.get("comment_rate","N/A"), "Comment Rate (per view)", size=32, dark_bg=dark)
        footer_bar(s, n, dark_bg=dark)

    # ── SLIDES 20-23: YT Shorts & Videos, each split in two ────
    yt_stats_slide(20, "YouTube — Shorts", "Shorts Performance", yt_shorts,
                   yt_raw.get("sample_size",30), "Shorts")
    deep_dive_slide(21, "YouTube — Shorts", "Shorts Deep Dive", yt_shorts,
                    T("youtube_shorts_analysis"), unit="Short")
    yt_stats_slide(22, "YouTube — Videos", "Videos Performance", yt_videos,
                   yt_raw.get("sample_size",30), "long-form videos")
    deep_dive_slide(23, "YouTube — Videos", "Videos Deep Dive", yt_videos,
                    T("youtube_videos_analysis"), unit="video")

    # ── SLIDE 24: YT Shorts vs Videos chart (light) ────────────
    s, dark = start_slide(prs, blank, 24)
    kicker_header(s, "YouTube", "Shorts vs Videos", "Comparing the two formats head-to-head", dark_bg=dark)
    def _views_num(v):
        try: return int(str(v).replace(",",""))
        except: return 0
    cd = ChartData(); cd.categories = ["Shorts", "Long-form Videos"]
    cd.add_series("Avg Views", (_views_num(yt_shorts.get("avg_views",0)), _views_num(yt_videos.get("avg_views",0))))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(1.80),
                                Inches(6.40), Inches(3.40), cd).chart
    chart.has_legend = False
    ser = chart.plots[0].series[0]
    for i in range(2):
        ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb = GOLD
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    tb(s, "AVG VIEWS PER UPLOAD", 0.55, 5.25, 6.00, 0.30, size=9.5,
       color=(TEXT_GRAY_LT if dark else TEXT_GRAY), font=FONT_MONO)
    stat_block(s, 7.45, 1.95, 5.00, yt_shorts.get("er_by_view","N/A"), "Shorts — ER by View", size=36, dark_bg=dark)
    stat_block(s, 7.45, 3.15, 5.00, yt_videos.get("er_by_view","N/A"), "Videos — ER by View", size=36, dark_bg=dark)
    sv = _views_num(yt_shorts.get("avg_views",0)); vv = _views_num(yt_videos.get("avg_views",0))
    ratio = f"{(vv/sv):.1f}x" if (sv and vv and vv >= sv) else (f"{(sv/vv):.1f}x" if (sv and vv) else "N/A")
    ratio_label = "Long-form View Advantage" if (sv and vv and vv >= sv) else "Shorts View Advantage"
    stat_block(s, 7.45, 4.35, 5.00, ratio, ratio_label, size=36, dark_bg=dark)
    card(s, 0.55, 5.70, 12.23, 1.20, "Analysis", T("youtube_comparison_analysis"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 24, dark_bg=dark)

    # ── SLIDE 25: LinkedIn Overview (dark) — no Employees ──────
    overview_slide(25, "LinkedIn", "LinkedIn Overview",
        f"{handles.get('linkedin','N/A')}  ·  Based on last {li_raw.get('sample_size','N/A')} posts",
        [(usfmt(li_raw.get("followers","N/A")), "Followers", 54),
         (str(li_raw.get("sample_size","N/A")), "Posts Analyzed", 54),
         (li_raw.get("posting_frequency","N/A"), "Posting Frequency", 44)],
        [(li_raw.get("engagement_rate","N/A"), "Engagement Rate / Follower", 32),
         (li_raw.get("avg_engagement","N/A"), "Avg Engagement (Total ÷ Posts)", 32),
         (li_raw.get("engagement_total","N/A"), "Total Engagement", 32)],
        li_raw, T("linkedin_analysis"))

    # ── SLIDE 26: LI Content Strategy (light) ──────────────────
    li_cat_order = [
        ("videos",       "Videos"),
        ("multi_image",  "Multi-Image"),
        ("single_image", "Single-Image"),
        ("documents",    "Documents"),
        ("polls",        "Polls"),
        ("text_only",    "Text-Only"),
        ("other",        "Others (Articles, Newsletters)"),
    ]
    s, dark = start_slide(prs, blank, 26)
    kicker_header(s, "LinkedIn", "Content Strategy",
                  f"Content mix from the last {li_raw.get('sample_size',30)} posts", dark_bg=dark)
    li_counts = [li_categories.get(key, {}).get("n", 0) for key, _ in li_cat_order]
    li_total  = max(sum(li_counts), 1)
    cd = ChartData(); cd.categories = [label for _, label in li_cat_order]
    cd.add_series("Posts", tuple(li_counts))
    chart = s.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.55),Inches(1.80),Inches(5.40),Inches(3.90),cd).chart
    chart.has_legend = True; chart.legend.font.size = Pt(9)
    LI_COLORS = [GOLD, RGBColor(0x8B,0x6F,0x47), CARD_DARK, PIE_TAN,
                 RGBColor(0x6B,0x59,0x40), RGBColor(0xA8,0xA2,0x96), RGBColor(0x4A,0x38,0x26)]
    for i, col in enumerate(LI_COLORS):
        pt = chart.plots[0].series[0].points[i]
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = col
    yy = 1.90
    for key, label in li_cat_order:
        cnt = li_categories.get(key, {}).get("n", 0)
        tb(s, label, 6.45, yy, 3.65, 0.30, size=12, bold=True,
           color=(TEXT_LIGHT if dark else TEXT_DARK), font=FONT_MONO_MED)
        tb(s, f"{int(round(cnt/li_total*100))}%  —  {cnt} post{'s' if cnt != 1 else ''}",
           10.15, yy, 2.60, 0.30, size=12, color=GOLD, font=FONT_MONO_MED)
        yy += 0.40
    card(s, 0.55, 5.90, 12.23, 1.05, "Content Strategy Analysis",
         T("linkedin_content_analysis", max_chars=300), dark_bg=dark, body_size=11)
    footer_bar(s, 26, dark_bg=dark)

    # ── SLIDES 27-33: LI per-format performance (analysis, no count stat) ──
    li_cat_analyses = A.get("linkedin_category_analyses", {}) or {}
    li_sample = li_raw.get("sample_size", 30)
    for idx, (key, label) in enumerate(li_cat_order):
        n = 27 + idx
        d = li_categories.get(key, {"n": 0})
        s, dark = start_slide(prs, blank, n)
        kicker_header(s, f"LinkedIn — {label}", f"{label} Performance",
                      f"Based on {d.get('n',0)} post{'s' if d.get('n',0) != 1 else ''} from the last {li_sample}",
                      dark_bg=dark)
        if d.get("n", 0) > 0:
            stat_block(s, 0.55, 1.75, 3.86, d.get("avg_engagement","N/A"), "Avg Engagement / Post", size=44, dark_bg=dark)
            stat_block(s, 4.73, 1.75, 3.86, d.get("er_per_follower","N/A"), "Engagement Rate / Follower", size=44, dark_bg=dark)
            stat_block(s, 8.92, 1.75, 3.86, d.get("avg_likes","N/A"), "Avg Likes / Post", size=32, dark_bg=dark)
            stat_block(s, 0.55, 3.10, 3.86, d.get("avg_comments","N/A"), "Avg Comments / Post", size=32, dark_bg=dark)
            stat_block(s, 4.73, 3.10, 3.86, d.get("avg_shares","N/A"), "Avg Reposts / Post", size=32, dark_bg=dark)
            link_list_card(s, 0.55, 4.20, 12.23, 0.85, "Top / Worst", [
                (f"Top (score {d.get('top_score','N/A')}):  ", d.get("top_url","N/A")),
                (f"Worst (score {d.get('worst_score','N/A')}):  ", d.get("worst_url","N/A")),
            ], dark_bg=dark)
        else:
            card(s, 0.55, 1.90, 12.23, 1.60, "No Posts in This Category",
                 f"No {label.lower()} posts appeared in the sample of LinkedIn posts analyzed for this report.",
                 dark_bg=dark)
        cat_analysis = truncate_to_sentence(str(li_cat_analyses.get(key, "") or ""), max_chars=420)
        card(s, 0.55, 5.30, 12.23, 1.55, "Analysis", cat_analysis, dark_bg=dark)
        footer_bar(s, n, dark_bg=dark)

    # ── SLIDE 34: LI Format Comparison (light) ─────────────────
    s, dark = start_slide(prs, blank, 34)
    kicker_header(s, "LinkedIn", "Format Comparison",
                  "Engagement rate per follower across all formats", dark_bg=dark)
    rows = []
    for key, label in li_cat_order:
        d = li_categories.get(key, {"n": 0})
        rows.append([label, d.get("n", 0), d.get("er_per_follower","—") if d.get("n",0) else "—",
                     d.get("avg_engagement","—") if d.get("n",0) else "—"])
    def _er_num(v):
        try: return float(str(v).replace("%",""))
        except: return -1
    rows.sort(key=lambda r: _er_num(r[2]), reverse=True)
    rounded_table(s, ["Format", "Posts", "ER / Followers", "Avg Engagement"], rows,
                  0.55, 1.72, 12.23, dark_bg=dark, row_h=0.44)
    card(s, 0.55, 6.05, 12.23, 1.00, "What Works Best",
         T("linkedin_comparison_analysis", max_chars=300), dark_bg=dark, body_size=11)
    footer_bar(s, 34, dark_bg=dark)

    # ── SLIDE 35: LI Summary (dark) ────────────────────────────
    summary_slide(35, "LinkedIn", "LinkedIn Summary", [
        ("Audience", f"{usfmt(li_raw.get('followers','N/A'))} followers · core B2B channel"),
        ("Performance", f"ER {li_raw.get('engagement_rate','N/A')} per follower · avg engagement {li_raw.get('avg_engagement','N/A')} / post · momentum {li_raw.get('momentum_pct','N/A')}"),
        ("Cadence", f"{li_raw.get('posting_frequency','N/A')} · {li_raw.get('posting_consistency','N/A')} · best day {li_raw.get('best_performing_day','N/A')}, {li_raw.get('best_performing_hour','N/A')}"),
        ("Strongest Format", str(li.get("strength","N/A"))),
        ("Overall LinkedIn Analysis", T("linkedin_summary")),
    ])

    # ── SLIDE 36: Cross-Platform Follower Comparison (light) ───
    s, dark = start_slide(prs, blank, 36)
    kicker_header(s, "Cross-Platform", "Follower Comparison", "Audience size across all platforms", dark_bg=dark)
    ig_f = parse_num(ig_raw.get("followers",0)); fb_f = parse_num(fb_raw.get("followers",0))
    yt_f = parse_num(yt_raw.get("subscribers",0)); li_f = parse_num(li_raw.get("followers",0))
    cd = ChartData(); cd.categories = ["Instagram","Facebook","YouTube","LinkedIn"]
    cd.add_series("Followers",(ig_f,fb_f,yt_f,li_f))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(0.55),Inches(1.75),Inches(12.23),Inches(3.55),cd).chart
    chart.has_legend = False
    ser = chart.plots[0].series[0]
    for i in range(4):
        ser.points[i].format.fill.solid(); ser.points[i].format.fill.fore_color.rgb = GOLD
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    for i,(lbl,val) in enumerate([
        ("Instagram", usfmt(ig_raw.get("followers","N/A"))),
        ("Facebook",  usfmt(fb_raw.get("followers","N/A"))),
        ("YouTube",   usfmt(yt_raw.get("subscribers","N/A"))),
        ("LinkedIn",  usfmt(li_raw.get("followers","N/A"))),
    ]):
        stat_block(s, 0.55 + i*3.13, 5.55, 2.83, val, lbl, size=32, dark_bg=dark)
    footer_bar(s, 36, dark_bg=dark)

    # ── SLIDE 37: All-Platform Summary (dark) ──────────────────
    s, dark = start_slide(prs, blank, 37)
    kicker_header(s, "All Platforms", "Platform Summary", "The four channels, side by side", dark_bg=dark)
    rounded_table(s, ["Platform", "Followers", "Engagement", "Frequency", "Momentum"], [
        ["Instagram", usfmt(ig_raw.get("followers","N/A")), ig_raw.get("engagement_rate","N/A"),
         ig_raw.get("posting_frequency","N/A"), ig_raw.get("momentum_pct","N/A")],
        ["Facebook",  usfmt(fb_raw.get("followers","N/A")), fb_raw.get("engagement_rate","N/A"),
         fb_raw.get("posting_frequency","N/A"), fb_raw.get("momentum_pct","N/A")],
        ["YouTube",   usfmt(yt_raw.get("subscribers","N/A")), yt_raw.get("engagement_rate","N/A"),
         yt_raw.get("posting_frequency","N/A"), yt_raw.get("momentum_pct","N/A")],
        ["LinkedIn",  usfmt(li_raw.get("followers","N/A")), li_raw.get("engagement_rate","N/A"),
         li_raw.get("posting_frequency","N/A"), li_raw.get("momentum_pct","N/A")],
    ], 0.55, 1.72, 12.23, dark_bg=dark)
    card(s, 0.55, 4.60, 12.23, 2.25, "Cross-Platform Summary", T("cross_platform_summary"), dark_bg=dark)
    footer_bar(s, 37, dark_bg=dark)

    # ── SLIDE 38: Strengths & Gaps (light) ─────────────────────
    s, dark = start_slide(prs, blank, 38)
    kicker_header(s, "Strengths & Gaps", "Key Strengths & Gaps",
                  "What's working, and what needs attention", dark_bg=dark)
    card(s, 0.55, 1.70, 5.96, 1.40, "Strongest Platform",
         truncate_to_sentence(str(cp.get("strongest_platform","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 6.83, 1.70, 5.96, 1.40, "Needs Most Work",
         truncate_to_sentence(str(cp.get("weakest_platform","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 0.55, 3.30, 5.96, 1.40, "Content Consistency",
         truncate_to_sentence(str(cp.get("content_consistency","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 6.83, 3.30, 5.96, 1.40, "Growth Opportunity",
         truncate_to_sentence(str(cp.get("growth_opportunity","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 0.55, 4.95, 12.23, 1.85, "Overall Analysis", T("overall_analysis"), dark_bg=dark, body_size=11.5)
    footer_bar(s, 38, dark_bg=dark)

    # ── SLIDE 39: Recommendations (dark) ───────────────────────
    s, dark = start_slide(prs, blank, 39)
    tb(s, "ACTION PLAN", 0.55, 0.42, 6.00, 0.30, size=11, color=GOLD, font=FONT_MONO_MED)
    tb(s, "Strategic Recommendations", 0.55, 0.72, 11.50, 0.65, size=26, color=TEXT_LIGHT, font=FONT_SERIF)
    tb(s, "Prioritised next steps by platform", 0.55, 1.34, 11.50, 0.35,
       size=13, color=TEXT_GRAY_LT, font=FONT_MONO_LT)
    card(s, 0.55, 1.75, 5.96, 1.45, "Instagram", truncate_to_sentence(str(ig.get("recommendation","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 6.83, 1.75, 5.96, 1.45, "YouTube",   truncate_to_sentence(str(yt.get("recommendation","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 0.55, 3.40, 5.96, 1.45, "Facebook",  truncate_to_sentence(str(fb.get("recommendation","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 6.83, 3.40, 5.96, 1.45, "LinkedIn",  truncate_to_sentence(str(li.get("recommendation","N/A")), 220), dark_bg=dark, body_size=11.5)
    card(s, 0.55, 5.10, 12.23, 1.30, "Top Priority",
         truncate_to_sentence(str(cp.get("overall_recommendation","N/A")), 420), dark_bg=dark)
    footer_bar(s, 39, dark_bg=dark)

    # ── SLIDE 40: Thank You (light) ────────────────────────────
    s, dark = start_slide(prs, blank, 40)
    tb(s, "Thank You", 0.55, 2.80, 12.23, 1.20, size=54,
       color=(TEXT_LIGHT if dark else TEXT_DARK), font=FONT_SERIF, align=PP_ALIGN.CENTER)
    box = s.shapes.add_textbox(Inches(0.55), Inches(4.30), Inches(12.23), Inches(0.40))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Generated by "
    r1.font.size = Pt(13); r1.font.name = FONT_MONO_LT
    r1.font.color.rgb = TEXT_GRAY_LT if dark else TEXT_GRAY
    r2 = p.add_run(); r2.text = "midnightpanda.ai"
    r2.hyperlink.address = "https://midnightpanda.ai"
    r2.font.size = Pt(13); r2.font.name = FONT_MONO_MED; r2.font.color.rgb = GOLD
    tb(s, "AI Brand Intelligence  ·  Competitor Ad Analysis  ·  Social Listening  ·  SEO & Traffic Intelligence",
       0.55, 4.80, 12.23, 0.35, size=11, color=(TEXT_GRAY_LT if dark else TEXT_GRAY),
       font=FONT_MONO_LT, align=PP_ALIGN.CENTER)
    footer_bar(s, 40, dark_bg=dark)

    # ── Save + force gold hyperlink theme ──────────────────────
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)
    safe_name = re.sub(r'[^A-Za-z0-9_]', '_', analysis.get("brand_name", website_url).replace(" ","_"))
    path = os.path.join(output_dir, f"{safe_name}_social_report.pptx")
    prs.save(path)
    try:
        _force_gold_hyperlinks(path)
        print("   ✅ Theme hyperlink color forced to gold")
    except Exception as e:
        print(f"   ⚠️ Gold hyperlink theme fix failed (links may render blue): {e}")
    print(f"✅ PPT saved: {path}")
    return path, safe_name


# ═══════════════════════════════════════════════════════════════
# ASYNC JOB SYSTEM
# ═══════════════════════════════════════════════════════════════
#
# WHY THIS EXISTS: the old /generate route did all the work — handle
# discovery, 4 platform fetches, Claude analysis, PPT building — inside
# one single HTTP request. As the pipeline grew (especially the new
# shares/reposts enrichment step), that request started taking long
# enough that the browser or Railway would give up waiting and show a
# generic "upstream error", EVEN THOUGH THE SERVER KEPT WORKING AND
# FINISHED THE REPORT ANYWAY IN THE BACKGROUND. Proof of this: the logs
# showed "PPT saved" and a 200 response for a request the user had
# already seen fail.
#
# THE FIX: /generate now only does one thing — kick off the real work
# on a background thread and return a job_id immediately (in well under
# a second, every time). The page then polls /status/<job_id> every
# couple of seconds until the job is done. No single request is ever
# open long enough to time out, no matter how slow a report gets.
#
# IMPORTANT DEPLOYMENT NOTE: job progress is stored in memory (JOBS
# dict below), not a database. This is intentional — simple and
# sufficient for a low-traffic internal tool — but it ONLY works
# correctly if the app runs as a SINGLE process/worker. If your
# Procfile or Railway start command runs gunicorn with more than one
# worker (e.g. "--workers 2"), a job started on worker A may not be
# visible when a later /status request lands on worker B, and polling
# would incorrectly show "not found". Check your Procfile — it should
# specify a single worker, e.g.:
#     web: gunicorn app:app --workers 1 --timeout 300
# A single worker is completely fine for one agency's internal usage.

JOBS = {}
JOBS_LOCK = threading.Lock()
JOB_MAX_AGE_SECONDS = 2 * 60 * 60  # auto-forget jobs older than 2 hours

def _cleanup_old_jobs():
    cutoff = time.time() - JOB_MAX_AGE_SECONDS
    stale = [jid for jid, j in JOBS.items() if j.get("created_at", 0) < cutoff]
    for jid in stale:
        del JOBS[jid]

def _job_update(job_id, **fields):
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id].update(fields)

def _job_add_step(job_id, label):
    with JOBS_LOCK:
        if job_id in JOBS:
            JOBS[job_id]["steps"].append(label)

def run_generate_job(job_id, website_url):
    """Runs the full report pipeline on a background thread. Never touches
    the HTTP request/response cycle directly — only writes progress into
    the JOBS dict, which /status/<job_id> reads."""
    try:
        _job_update(job_id, status="running")

        _job_add_step(job_id, "Scanning website for social handles")
        handles = discover_all_handles(website_url)

        ig_raw = {}
        if handles.get("instagram"):
            _job_add_step(job_id, "Fetching Instagram profile data")
            ig_raw = fetch_instagram(handles.get("instagram"))

        fb_raw = {}
        if handles.get("facebook"):
            _job_add_step(job_id, "Fetching Facebook page data")
            fb_raw = fetch_facebook(handles.get("facebook"))

        yt_raw = {}
        if handles.get("youtube"):
            _job_add_step(job_id, "Fetching YouTube channel data")
            yt_raw = fetch_youtube(handles.get("youtube"))

        domain_for_brand = website_url.replace("https://","").replace("http://","").rstrip("/")
        brand_for_search = re.sub(r'\.(com|in|io|co|net|org|app).*','', domain_for_brand).replace("www.","")
        li_raw = {}
        if handles.get("linkedin"):
            _job_add_step(job_id, "Fetching LinkedIn data")
            li_raw = fetch_linkedin(handles.get("linkedin"), brand_for_search)

        _job_add_step(job_id, "Analysing with Claude AI")
        analysis = analyse_with_claude(website_url, handles, ig_raw, fb_raw, yt_raw, li_raw)

        _job_add_step(job_id, f"Building {TOTAL_SLIDES}-slide PowerPoint")
        ppt_path, safe_name = create_ppt(analysis, handles, ig_raw, fb_raw, yt_raw, li_raw, website_url)

        _job_update(job_id, status="done", result={
            "success":      True,
            "brand":        analysis.get("brand_name",""),
            "handles":      handles,
            "analysis":     analysis,
            "download_url": f"/download/{safe_name}"
        })

    except ValueError as e:
        _job_update(job_id, status="error", error=str(e), error_code=404)
    except Exception as e:
        import traceback
        print(traceback.format_exc())
        _job_update(job_id, status="error", error=str(e), error_code=500)


# ═══════════════════════════════════════════════════════════════
# FLASK ROUTES
# ═══════════════════════════════════════════════════════════════

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    body        = request.get_json()
    website_url = (body or {}).get("website_url","").strip()
    if not website_url:
        return jsonify({"error":"Please enter a website URL"}), 400

    with JOBS_LOCK:
        _cleanup_old_jobs()
        job_id = str(uuid.uuid4())
        JOBS[job_id] = {
            "status": "starting",
            "steps": [],
            "result": None,
            "error": None,
            "created_at": time.time(),
        }

    thread = threading.Thread(target=run_generate_job, args=(job_id, website_url), daemon=True)
    thread.start()

    # Returns almost instantly — this request never does the slow work,
    # so it can never be the thing that times out.
    return jsonify({"job_id": job_id}), 202

@app.route("/status/<job_id>")
def job_status(job_id):
    with JOBS_LOCK:
        job = JOBS.get(job_id)
        if not job:
            return jsonify({"error": "Job not found (it may have expired, or the server restarted)"}), 404
        snapshot = dict(job)  # copy fields out before releasing the lock

    if snapshot["status"] == "error":
        return jsonify({
            "status": "error",
            "steps":  snapshot["steps"],
            "error":  snapshot["error"],
        }), snapshot.get("error_code", 500)

    if snapshot["status"] == "done":
        return jsonify({
            "status": "done",
            "steps":  snapshot["steps"],
            **snapshot["result"],
        })

    return jsonify({"status": snapshot["status"], "steps": snapshot["steps"]})

@app.route("/download/<safe_name>")
def download(safe_name):
    path = f"output/{safe_name}_social_report.pptx"
    if not os.path.exists(path):
        return jsonify({"error":"Report not found"}), 404
    return send_file(path, as_attachment=True,
                     download_name=f"{safe_name}_social_report.pptx")

@app.route("/debug_repost")
def debug_repost():
    """Quick test of the analytics actor for one post URL — returns the
    shares/reposts/saves it found, for verifying the enrichment source."""
    post_url = request.args.get("url", "")
    if not post_url:
        return jsonify({"error": "Pass a post URL like ?url=https://www.instagram.com/p/XXXX/"})
    results = fetch_post_metrics_via_analytics_actor([post_url])
    return jsonify({
        "requested_url": post_url,
        "metrics_found": results.get(post_url) or results,
    })

@app.route("/debug_youtube/<channel_handle>")
def debug_youtube(channel_handle):
    """Run this to verify the official YouTube API is fetching real
    videos for a given channel. Unlike Instagram/Facebook/LinkedIn,
    field names here are guaranteed correct (official Google docs) —
    this route mainly confirms channel handle resolution works."""
    try:
        videos, resolved_handle = fetch_youtube_videos_via_api(channel_handle)
        return jsonify({
            "total_videos_returned": len(videos),
            "resolved_handle":       resolved_handle,
            "first_video_full":      videos[0] if videos else "no videos found",
            "second_video_full":     videos[1] if len(videos) > 1 else "n/a",
        })
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/debug_linkedin/<company_handle>")
def debug_linkedin(company_handle):
    """Run this once against a real LinkedIn company page to check the
    raw Apify output — field names in fetch_linkedin_posts_via_apify /
    _classify_linkedin_post are best-effort guesses and may need
    adjusting once you see real data, same process used for Facebook."""
    try:
        posts = fetch_linkedin_posts_via_apify(company_handle)
        return jsonify({
            "total_posts_returned": len(posts),
            "first_post_full":      posts[0] if posts else "no posts found",
            "second_post_full":     posts[1] if len(posts) > 1 else "n/a",
        })
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/debug_facebook/<username>")
def debug_facebook(username):
    """Run this once against a real Facebook Page to check the raw Apify
    output — the field names in fetch_facebook_posts_via_apify /
    _classify_facebook_post are best-effort guesses and may need
    adjusting once you see real data, same process used to verify
    Instagram's fields."""
    try:
        posts = fetch_facebook_posts_via_apify(username)
        return jsonify({
            "total_posts_returned": len(posts),
            "first_post_full":      posts[0] if posts else "no posts found",
            "second_post_full":     posts[1] if len(posts) > 1 else "n/a",
        })
    except Exception as e:
        return jsonify({"error": str(e)})

@app.route("/debug_instagram/<username>")
def debug_instagram(username):
    try:
        r = requests.get(
            "https://www.searchapi.io/api/v1/search",
            params={"engine": "instagram_profile", "username": username,
                    "api_key": SEARCHAPI_KEY}
        )
        data  = r.json()
        posts = data.get("posts", [])
        return jsonify({
            "status_code":          r.status_code,
            "total_posts_returned": len(posts),
            "search_metadata":      data.get("search_metadata", {}),
            "profile_keys":         list(data.get("profile", {}).keys()),
            "first_post_full":      posts[0] if posts else "no posts found",
            "second_post_full":     posts[1] if len(posts) > 1 else "n/a",
        })
    except Exception as e:
        return jsonify({"error": str(e)})

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5001))
    # threaded=True is required here: without it, Flask's built-in server
    # handles one request at a time, which would block /status polling
    # while a report is being generated in the background.
    app.run(host="0.0.0.0", port=port, debug=False, threaded=True)